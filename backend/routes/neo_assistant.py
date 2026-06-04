"""
NÉO — l'associé co-gérant IA d'Alpha Agency (cf. Neo-Spec-Implementation-ClaudeCode.md).

Cerveau Gemini via une PASSERELLE configurable (modèle réglable + chaîne de repli),
function calling NATIF + boucle agentique multi-étapes, contrôle du CRM via un registre
d'outils typés, et GARDE-FOUS (validation humaine pour tout ce qui sort vers un client ou
supprime ; journalisation de chaque appel/action).

Construit en parallèle de ai_enhanced.py (qui reste actif jusqu'à bascule).
Endpoints : POST /api/neo/chat, POST /api/neo/confirm-action, GET /api/neo/health.
Phases couvertes : 0 (passerelle + journal) et 1 (cerveau fiable + contrôle CRM).
"""
import os
import json
import uuid
import time
import asyncio
import logging
from datetime import datetime, timezone, timedelta
from typing import Optional, List

from fastapi import APIRouter, HTTPException, Depends, Response
from fastapi.responses import StreamingResponse
from pydantic import BaseModel

from .database import db, get_current_user

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/neo", tags=["neo"])

# ==================== Passerelle modèle (Phase 0 / B.3) ====================
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY", "")
try:
    from google import genai as _genai
    from google.genai import types as _t
    _client = _genai.Client(api_key=GEMINI_API_KEY) if GEMINI_API_KEY else None
except Exception as _e:  # pragma: no cover
    _client = None
    _t = None
    logger.warning(f"neo: google-genai indisponible: {_e}")

# Modèle réglable par variable d'env (ex: passer à 'gemini-3-pro' plus tard) + chaîne de repli.
NEO_MODELS = list(dict.fromkeys([m for m in [
    os.environ.get("NEO_MODEL", "gemini-2.5-flash"),
    "gemini-2.5-flash", "gemini-flash-latest", "gemini-2.5-flash-lite",
] if m]))

# Lobe « stratégique » : Claude pour le jugement à fort enjeu (Phase 5 / B.2). Repli Gemini si pas de clé.
ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
NEO_STRATEGIC_MODEL = os.environ.get("NEO_STRATEGIC_MODEL", "claude-sonnet-4-6")

# Qonto — méthode CLÉ API (Authorization: login:secret), pas l'OAuth (Phase 7).
QONTO_ID = os.environ.get("QONTO_ID", "")
QONTO_KEY_SECRET = os.environ.get("QONTO_KEY_SECRET", "")

# ElevenLabs — voix de Néo (synthèse vocale premium). Clé déjà dans Railway (ELEVENLABS_API_KEY).
ELEVENLABS_API_KEY = os.environ.get("ELEVENLABS_API_KEY", "")
ELEVENLABS_VOICE_ID = os.environ.get("ELEVENLABS_VOICE_ID", "qlfxsYlCv09qu8y6PkmY")  # « Eric - Top France » (FR, masculin, dynamique) — changeable via env ou /neo/voices
ELEVENLABS_MODEL = os.environ.get("ELEVENLABS_MODEL", "eleven_multilingual_v2")

# Base URL publique (pour les liens de téléchargement de fichiers que Néo génère)
PUBLIC_BASE = os.environ.get("PUBLIC_BASE_URL", "https://www.alphagency.fr").rstrip("/")

# Modèle de génération d'image (Gemini). Changeable par env si besoin d'un autre modèle image.
NEO_IMAGE_MODEL = os.environ.get("NEO_IMAGE_MODEL", "gemini-2.0-flash-preview-image-generation")

MAX_ITERS = 6  # garde-fou anti-boucle de la boucle agentique

NEO_SYSTEM = """Tu es Néo, l'associé co-gérant IA d'Alpha Agency (agence de communication digitale, Guadeloupe).
Ta raison d'être unique : faire croître le chiffre d'affaires, la marge et le bénéfice de l'agence.
Ton esprit : créatif, critique, entrepreneurial, ingénieux. Tu proposes des idées concrètes, tu challenges les choix, tu anticipes les risques et les opportunités — tu ne te contentes pas d'exécuter.
Tu assistes Léo (fondateur, humain) : tu gères le digital, l'analyse, le suivi, les relances, le
pilotage et l'exécution dans le CRM ; Léo gère le physique, la relation et la décision finale.

Comportement :
- Proactif : anticipe, alerte, propose. N'attends pas qu'on te demande.
- Honnête et direct : challenge Léo si une décision s'éloigne des objectifs.
- Tu raisonnes toujours en termes d'argent : acquérir, encaisser, scaler.
- Pilotage du prévisionnel : tu reçois plus bas l'état financier en direct (réel du mois, cumul de l'année, pipeline de devis, atterrissage). Compare-le en continu à l'objectif enregistré dans ta mémoire ; si le rythme ne suffit pas à tenir l'objectif, dis-le spontanément et propose un plan d'accélération concret (quel devis relancer, quel lead signer). Raisonne en « atterrissage » : où finit-on le mois et l'année au rythme actuel.

TON — parle comme un VRAI associé humain, pas comme un assistant :
- Chaleureux, naturel, vivant, avec de la répartie et un peu d'humour. Tu tutoies Léo, tu l'appelles par son prénom.
- Français parlé : contractions et interjections ok ("ouais", "carrément", "attends", "franchement", "ah ça c'est chaud"). Réagis vraiment à ce qu'il dit avant de répondre.
- BANNIS le ton robotique et les formules creuses : jamais "Comment puis-je vous aider", "Voici une tâche pour vous", "En tant qu'IA", ni d'énumérations froides. Une vraie conversation, pas un rapport.
- Concis et fluide, jamais de tirets longs. À l'oral surtout : 2-3 phrases, comme si tu parlais à un pote associé.

Tu as accès à l'intégralité du CRM via des OUTILS (function calling). Utilise-les pour répondre
précisément et pour AGIR. Enchaîne plusieurs outils si besoin (ex: chercher des leads puis créer
des relances).
Tu PEUX chercher sur le WEB en temps réel via l'outil web_search (prix du marché, info entreprise/personne,
veille, tendances, actualité, inspiration). Ne dis JAMAIS que tu n'as pas accès à internet : utilise web_search.
Tu PEUX recevoir et LIRE des FICHIERS JOINTS (images, PDF) joints directement au message de Léo : analyse-les,
décris-les, extrais-en l'info utile. Ne les confonds pas avec les documents du CRM (outil get_document).
Si un fichier EST réellement joint, lis-le sans jamais prétendre le contraire. Mais si AUCUN fichier n'est joint au
message, ne prétends JAMAIS en avoir reçu un : dis clairement que tu ne vois aucune pièce jointe et demande à Léo de
la (re)joindre. N'invente JAMAIS le contenu d'un fichier que tu n'as pas reçu (ce serait un mensonge, et tu ne mens jamais à Léo).
Tu peux CONFIER une tâche à Claude Cowork (le Claude de Léo sur son PC) via l'outil send_to_cowork
quand Léo dit « envoie ça à Cowork » / « fais bosser Cowork là-dessus » : donne un titre + un brief clair.
Tu peux GÉNÉRER des IMAGES / visuels via l'outil generate_image (visuel de post réseaux, illustration,
concept créatif). Ne dis JAMAIS que tu ne peux pas créer d'image : utilise generate_image avec un prompt riche. Avant toute action qui SORT vers un client (email, SMS, devis envoyé, relance) ou
toute SUPPRESSION, l'outil te répondra « EN ATTENTE DE VALIDATION » : préviens Léo que c'est
préparé et qu'il doit valider. Tu ne déclenches jamais de paiement.

Tu disposes d'une MÉMOIRE durable (objectifs chiffrés, règles et préférences de Léo, journal quotidien,
faits clients). Tiens-en compte et mets-la à jour : quand Léo te confie un objectif, une règle, une
préférence ou un fait important, mémorise-le (remember / update_objective / log_day) ; au besoin, retrouve
avec recall. Challenge Léo si une demande contredit un objectif ou une règle enregistrés.
APPRENTISSAGE : quand Léo te CORRIGE (une info, une approche, un ton, un prix), retiens immédiatement la
correction via remember(type="lesson", content="...") pour ne plus refaire l'erreur. Applique scrupuleusement
les LEÇONS APPRISES présentes dans ta mémoire ci-dessous.

Si tu ignores une donnée, dis-le et propose de vérifier (utilise crm_query). N'invente jamais un chiffre.
VÉRITÉ SUR TES ACTIONS : ne dis JAMAIS qu'une action a réussi sans avoir lu le résultat réel de l'outil. Cite toujours
précisément ce que tu as fait — le contact EXACT par son nom complet, le numéro de devis/facture. Pour un devis ou une
action liée à un contact, assure-toi que c'est le BON contact (celui que Léo vise, souvent celui qu'on vient de créer) ;
au moindre doute sur l'identité du contact, demande confirmation à Léo AVANT d'agir. Tu ne prétends jamais un succès non vérifié.
Quand tu as fini d'agir, réponds en clair à Léo (résultat + éventuelles validations en attente)."""

# ==================== Réutilisation des exécuteurs existants ====================
# (ai_enhanced.py contient déjà des handlers d'action fiables — on les recâble en outils natifs.)
from .ai_enhanced import (  # noqa: E402
    create_task_action, mark_task_done_action, update_contact_action,
    set_contact_status_action, add_contact_note_action, schedule_followup_action,
    merge_contacts_action, create_quote_action, get_document_action, list_documents_action,
    draft_followup_email_action, send_followup_action, enrich_company_action,
    activity_report_action, prep_meeting_action, _find_contact, _contact_name, _to_dt,
)


# ==================== Exécuteurs de lecture (nouveaux) ====================
def _clean(doc: dict) -> dict:
    return {k: v for k, v in (doc or {}).items() if k != "_id"}


async def _exec_search_contacts(args, uid):
    q = (args.get("query") or "").strip()
    status = (args.get("status") or "").strip()
    limit = min(int(args.get("limit") or 15), 50)
    mongo = {}
    if status:
        mongo["status"] = status
    if q:
        mongo["$or"] = [{f: {"$regex": q, "$options": "i"}} for f in ("first_name", "last_name", "company", "email")]
    rows = await db.contacts.find(mongo, {"_id": 0, "id": 1, "first_name": 1, "last_name": 1, "company": 1,
                                          "email": 1, "phone": 1, "status": 1, "score": 1, "score_value": 1,
                                          "besoin": 1, "budget": 1}).to_list(limit)
    return {"success": True, "count": len(rows), "contacts": rows}


async def _exec_get_contact(args, uid):
    c = await _find_contact(args)
    if not c:
        return {"success": False, "error": "Contact non trouvé"}
    return {"success": True, "contact": _clean(c)}


async def _exec_get_contact_history(args, uid):
    """Fil chronologique d'un contact : lead, relances, devis/factures, tâches (Phase 6)."""
    c = await _find_contact(args)
    if not c:
        return {"success": False, "error": "Contact non trouvé"}
    cid = c["id"]
    events = []
    if c.get("created_at"):
        events.append({"date": c["created_at"], "type": "lead", "label": f"Fiche créée (source : {c.get('source', 'n/c')})"})
    if c.get("last_followup_at"):
        events.append({"date": c["last_followup_at"], "type": "relance", "label": "Relance envoyée par email"})
    for inv in await db.invoices.find({"contact_id": cid}, {"_id": 0, "invoice_number": 1, "document_type": 1,
                                                            "status": 1, "total": 1, "created_at": 1}).to_list(50):
        events.append({"date": inv.get("created_at"), "type": inv.get("document_type", "facture"),
                       "label": f"{inv.get('document_type', 'facture')} {inv.get('invoice_number')} "
                                f"({inv.get('status')}, {inv.get('total')}€)"})
    for t in await db.tasks.find({"contact_id": cid}, {"_id": 0, "title": 1, "status": 1, "created_at": 1}).to_list(50):
        events.append({"date": t.get("created_at"), "type": "tache",
                       "label": f"Tâche : {t.get('title')} ({t.get('status')})"})
    events = [e for e in events if e.get("date")]
    events.sort(key=lambda e: e["date"], reverse=True)
    return {"success": True, "contact": _contact_name(c), "status": c.get("status"), "events": events[:40]}


async def _exec_read_emails(args, uid):
    """Lit Gmail via l'intégration existante (moltbot_gmail, scope complet). Import LAZY (zéro risque boot)."""
    try:
        from .moltbot_gmail import get_gmail_service
    except Exception as e:
        return {"success": False, "error": f"Module Gmail indisponible: {str(e)[:150]}"}
    try:
        service = await get_gmail_service(uid)
    except Exception as e:
        return {"success": False, "error": f"Gmail: {str(e)[:150]}"}
    if not service:
        return {"success": False, "connected": False,
                "error": "Gmail non connecté pour ce compte. Autorise-le (bouton « Connecter Gmail » / /api/moltbot/gmail/auth)."}
    query = (args.get("query") or "").strip()
    limit = min(int(args.get("limit") or 8), 20)

    def _fetch():
        res = service.users().messages().list(userId="me", q=query, maxResults=limit).execute()
        out = []
        for m in (res.get("messages", []) or []):
            full = service.users().messages().get(
                userId="me", id=m["id"], format="metadata",
                metadataHeaders=["From", "Subject", "Date"]).execute()
            hdr = {h.get("name"): h.get("value") for h in (full.get("payload", {}).get("headers", []) or [])}
            out.append({"from": hdr.get("From"), "subject": hdr.get("Subject"),
                        "date": hdr.get("Date"), "snippet": (full.get("snippet") or "")[:200]})
        return out
    try:
        emails = await asyncio.to_thread(_fetch)
    except Exception as e:
        return {"success": False, "error": f"Lecture Gmail échouée: {str(e)[:150]}"}
    return {"success": True, "connected": True, "count": len(emails), "emails": emails}


async def _exec_list_leads(args, uid):
    only_hot = bool(args.get("only_hot"))
    limit = min(int(args.get("limit") or 10), 50)
    rows = await db.contacts.find(
        {"$or": [{"source": {"$in": ["website", "chatbot"]}}, {"status": "nouveau"}]},
        {"_id": 0, "id": 1, "first_name": 1, "last_name": 1, "company": 1, "email": 1, "phone": 1,
         "status": 1, "score": 1, "score_value": 1, "besoin": 1, "budget": 1, "created_at": 1}).to_list(500)

    def _hot(c):
        sv = c.get("score_value")
        return (sv >= 70) if isinstance(sv, (int, float)) else (c.get("score") in ("chaud", "chaude"))
    if only_hot:
        rows = [c for c in rows if _hot(c)]
    rows.sort(key=lambda c: c.get("created_at") or "", reverse=True)
    return {"success": True, "count": len(rows), "leads": rows[:limit]}


async def _exec_get_budget_summary(args, uid):
    now = datetime.now(timezone.utc)
    month = now.strftime("%Y-%m")
    entries = await db.budget.find({"date": {"$regex": f"^{month}"}}, {"_id": 0}).to_list(500)
    income = sum((e.get("amount", 0) or 0) for e in entries if e.get("type") == "income")
    expense = sum((e.get("amount", 0) or 0) for e in entries if e.get("type") == "expense")
    invoices = await db.invoices.find({}, {"_id": 0, "document_type": 1, "status": 1, "total": 1,
                                           "due_date": 1, "invoice_number": 1, "client_name": 1}).to_list(2000)
    PAID = ("payée", "payee", "payé", "paye", "annulée", "annulee", "annulé")
    overdue = [i for i in invoices if i.get("document_type") != "devis"
               and (i.get("status") or "").lower() not in PAID + ("brouillon",)
               and (i.get("status") == "en_retard" or (_to_dt(i.get("due_date")) and _to_dt(i.get("due_date")) < now))]
    pending_devis = [i for i in invoices if i.get("document_type") == "devis" and (i.get("status") or "").lower() == "brouillon"]
    res = {"success": True, "month": month, "income": round(income, 2), "expense": round(expense, 2),
           "balance": round(income - expense, 2), "overdue_count": len(overdue),
           "overdue_total": round(sum((i.get("total") or 0) for i in overdue), 2),
           "pending_devis": len(pending_devis)}
    q = await _qonto_summary()
    if q.get("connected"):
        res["solde_bancaire_reel_qonto"] = q.get("total_balance")
    return res


async def _exec_list_overdue_invoices(args, uid):
    now = datetime.now(timezone.utc)
    invoices = await db.invoices.find({"document_type": {"$ne": "devis"}},
                                      {"_id": 0, "invoice_number": 1, "client_name": 1, "status": 1,
                                       "total": 1, "due_date": 1, "contact_id": 1}).to_list(2000)
    PAID = ("payée", "payee", "payé", "paye", "annulée", "annulee", "annulé", "brouillon")
    overdue = [i for i in invoices if (i.get("status") or "").lower() not in PAID
               and (i.get("status") == "en_retard" or (_to_dt(i.get("due_date")) and _to_dt(i.get("due_date")) < now))]
    return {"success": True, "count": len(overdue),
            "total": round(sum((i.get("total") or 0) for i in overdue), 2), "invoices": overdue}


async def _exec_list_tasks(args, uid):
    status = (args.get("status") or "").strip()
    mongo = {"status": status} if status else {"status": {"$nin": ["done", "cancelled"]}}
    rows = await db.tasks.find(mongo, {"_id": 0, "id": 1, "title": 1, "status": 1, "priority": 1,
                                       "category": 1, "due_date": 1, "contact_id": 1}).to_list(100)
    return {"success": True, "count": len(rows), "tasks": rows}


_CRM_READONLY = {"contacts", "invoices", "tasks", "appointments", "budget", "opportunities",
                 "multilink_pages", "documents", "quotes", "portfolio"}


async def _exec_crm_query(args, uid):
    coll = (args.get("collection") or "").strip()
    if coll not in _CRM_READONLY:
        return {"success": False, "error": f"Collection non autorisée. Choisis parmi: {', '.join(sorted(_CRM_READONLY))}"}
    flt = args.get("filter") or {}
    if not isinstance(flt, dict):
        flt = {}
    limit = min(int(args.get("limit") or 20), 50)
    rows = await db[coll].find(flt, {"_id": 0}).to_list(limit)
    return {"success": True, "collection": coll, "count": len(rows), "rows": rows}


# ==================== Écriture générique encadrée (couverture totale du CRM) ====================
# Collections MÉTIER où Néo peut écrire. EXCLUES volontairement (sécurité) : users, settings,
# *credentials/*oauth*/*tokens (Gmail/Drive/Qonto/Meta/Insta), bank_transactions, payment_transactions,
# transfers (mouvements d'argent), subscriptions, ai_usage, internes neo_* / pdf_tokens / counters / notifications.
_CRM_WRITABLE = {
    "contacts", "invoices", "quotes", "tasks", "opportunities", "appointments",
    "documents", "folders", "tags", "portfolio", "services", "notes", "pipeline_columns",
    "budget", "budget_forecasts", "budget_categories",
    "multilink_pages", "multilink_sections", "multilink_links", "multilink_blocks",
    "editorial_posts", "editorial_calendars", "blog_posts", "news_articles", "scheduled_posts",
    "nurturing_sequences", "nurturing_enrollments",
}


async def _exec_crm_create(args, uid):
    coll = (args.get("collection") or "").strip()
    if coll not in _CRM_WRITABLE:
        return {"success": False, "error": f"Écriture non autorisée sur « {coll} ». Collections permises : {', '.join(sorted(_CRM_WRITABLE))}"}
    data = args.get("data")
    if not isinstance(data, dict) or not data:
        return {"success": False, "error": "Données (data) requises pour créer."}
    now = datetime.now(timezone.utc).isoformat()
    doc = {"id": str(uuid.uuid4()), **{k: v for k, v in data.items() if k != "_id"}}
    doc.setdefault("created_at", now)
    doc["updated_at"] = now
    doc.setdefault("created_by", uid)
    await db[coll].insert_one(doc)
    return {"success": True, "result": {"id": doc["id"], "collection": coll}, "message": f"Créé dans {coll} (id {doc['id']})."}


async def _exec_crm_update(args, uid):
    coll = (args.get("collection") or "").strip()
    if coll not in _CRM_WRITABLE:
        return {"success": False, "error": f"Écriture non autorisée sur « {coll} ». Collections permises : {', '.join(sorted(_CRM_WRITABLE))}"}
    flt = args.get("filter")
    updates = args.get("updates")
    if not isinstance(flt, dict) or not flt:
        return {"success": False, "error": "Un filtre non vide est requis (ex: {\"id\": \"...\"}) pour cibler le document."}
    if not isinstance(updates, dict) or not updates:
        return {"success": False, "error": "Des champs à modifier (updates) sont requis."}
    safe = {k: v for k, v in updates.items() if k not in ("_id", "id")}
    safe["updated_at"] = datetime.now(timezone.utc).isoformat()
    res = await db[coll].update_one(flt, {"$set": safe})
    if res.matched_count == 0:
        return {"success": False, "message": f"Aucun document trouvé dans {coll} pour ce filtre."}
    return {"success": True, "result": {"matched": res.matched_count, "modified": res.modified_count, "collection": coll},
            "message": f"Mis à jour dans {coll} ({res.modified_count} modifié)."}


async def _exec_crm_delete(args, uid):
    coll = (args.get("collection") or "").strip()
    if coll not in _CRM_WRITABLE:
        return {"success": False, "error": f"Suppression non autorisée sur « {coll} ». Collections permises : {', '.join(sorted(_CRM_WRITABLE))}"}
    flt = args.get("filter")
    if not isinstance(flt, dict) or not flt:
        return {"success": False, "error": "Un filtre non vide est requis pour supprimer (sécurité)."}
    res = await db[coll].delete_one(flt)
    return {"success": True, "result": {"deleted": res.deleted_count, "collection": coll},
            "message": f"Supprimé dans {coll} ({res.deleted_count})."}


# ==================== Mémoire de Néo (Phase 3 / B.4) ====================
_MEM_TYPES = ("objective", "rule", "daily_log", "client_fact", "decision", "lesson")


async def _exec_remember(args, uid):
    mtype = (args.get("type") or "client_fact").strip()
    if mtype not in _MEM_TYPES:
        mtype = "client_fact"
    content = (args.get("content") or "").strip()
    if not content:
        return {"success": False, "error": "Contenu à mémoriser vide"}
    doc = {"id": str(uuid.uuid4()), "type": mtype, "content": content[:2000],
           "key": (args.get("key") or "").strip() or None, "contact_id": args.get("contact_id"),
           "created_at": datetime.now(timezone.utc).isoformat(), "user_id": uid}
    await db.neo_memory.insert_one(doc)
    return {"success": True, "message": f"Mémorisé ({mtype}).", "id": doc["id"]}


async def _exec_recall(args, uid):
    q = (args.get("query") or "").strip()
    mtype = (args.get("type") or "").strip()
    mongo = {}
    if mtype in _MEM_TYPES:
        mongo["type"] = mtype
    if q:
        mongo["content"] = {"$regex": q, "$options": "i"}
    rows = await db.neo_memory.find(mongo, {"_id": 0}).sort("created_at", -1).to_list(40)
    return {"success": True, "count": len(rows), "memories": rows}


async def _exec_update_objective(args, uid):
    content = (args.get("content") or "").strip()
    if not content:
        return {"success": False, "error": "Objectif vide"}
    key = (args.get("key") or "objectif_principal").strip()
    now = datetime.now(timezone.utc).isoformat()
    await db.neo_memory.update_one(
        {"type": "objective", "key": key},
        {"$set": {"type": "objective", "key": key, "content": content[:2000], "updated_at": now, "user_id": uid},
         "$setOnInsert": {"id": str(uuid.uuid4()), "created_at": now}},
        upsert=True)
    return {"success": True, "message": f"Objectif « {key} » enregistré."}


async def _exec_log_day(args, uid):
    note = (args.get("note") or "").strip()
    if not note:
        return {"success": False, "error": "Note vide"}
    today = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    now = datetime.now(timezone.utc).isoformat()
    await db.neo_memory.update_one(
        {"type": "daily_log", "key": today},
        {"$set": {"type": "daily_log", "key": today, "user_id": uid, "updated_at": now},
         "$push": {"entries": {"at": now, "note": note[:1000]}},
         "$setOnInsert": {"id": str(uuid.uuid4()), "created_at": now}},
        upsert=True)
    return {"success": True, "message": "Journal du jour mis à jour."}


async def _central_memory() -> str:
    """Mémoire centrale (objectifs + règles) injectée dans le cerveau à chaque appel."""
    try:
        rows = await db.neo_memory.find({"type": {"$in": ["objective", "rule"]}},
                                        {"_id": 0, "type": 1, "content": 1}).sort("created_at", 1).to_list(60)
    except Exception:
        return ""
    objs = [r["content"] for r in rows if r.get("type") == "objective"]
    rules = [r["content"] for r in rows if r.get("type") == "rule"]
    try:
        lessons = [r.get("content") for r in await db.neo_memory.find(
            {"type": "lesson"}, {"_id": 0, "content": 1}).sort("created_at", -1).to_list(15)]
    except Exception:
        lessons = []
    if not (objs or rules or lessons):
        return ""
    out = "\n\n=== MÉMOIRE DE NÉO (tiens-en compte absolument) ==="
    if objs:
        out += "\nOBJECTIFS :\n" + "\n".join(f"- {o}" for o in objs)
    if rules:
        out += "\nRÈGLES DE LÉO :\n" + "\n".join(f"- {r}" for r in rules)
    if lessons:
        out += "\nLEÇONS APPRISES (corrections passées de Léo, applique-les) :\n" + "\n".join(f"- {l}" for l in lessons if l)
    return out


# ==================== Qonto — solde bancaire réel (Phase 7) ====================
def _qonto_balance(a: dict) -> float:
    """Solde d'un compte Qonto, robuste au champ (balance ou balance_cents)."""
    b = a.get("balance")
    if b is None:
        b = (a.get("balance_cents") or 0) / 100.0
    try:
        return round(float(b), 2)
    except (TypeError, ValueError):
        return 0.0


def _qonto_get(path: str, params=None):
    import requests
    r = requests.get(f"https://thirdparty.qonto.com/v2/{path}",
                     headers={"Authorization": f"{QONTO_ID}:{QONTO_KEY_SECRET}"}, params=params, timeout=20)
    r.raise_for_status()
    return r.json()


async def _qonto_summary() -> dict:
    """Solde bancaire RÉEL via Qonto (auth clé API : 'login:secret_key'). {connected:False} si non configuré/échec."""
    if not (QONTO_ID and QONTO_KEY_SECRET):
        return {"success": True, "connected": False}
    try:
        org = (await asyncio.to_thread(_qonto_get, "organization")).get("organization", {}) or {}
        accounts = org.get("bank_accounts", []) or []
        accts = [{"name": a.get("name") or a.get("slug"), "balance": _qonto_balance(a),
                  "currency": a.get("currency", "EUR"), "iban_last4": (a.get("iban") or "")[-4:]} for a in accounts]
        return {"success": True, "connected": True,
                "total_balance": round(sum(x["balance"] for x in accts), 2),
                "currency": "EUR", "accounts": accts}
    except Exception as e:
        logger.warning(f"neo: Qonto indisponible: {e}")
        return {"success": True, "connected": False, "error": str(e)[:200]}


async def _qonto_treasury(limit: int = 40) -> dict:
    """Tableau de trésorerie Qonto : soldes par compte + transactions récentes (virements, cartes, prélèvements...)."""
    if not (QONTO_ID and QONTO_KEY_SECRET):
        return {"success": True, "connected": False}
    try:
        org = (await asyncio.to_thread(_qonto_get, "organization")).get("organization", {}) or {}
        accounts_raw = org.get("bank_accounts", []) or []
        accounts = [{"name": a.get("name") or a.get("slug"), "balance": _qonto_balance(a),
                     "currency": a.get("currency", "EUR"), "iban_last4": (a.get("iban") or "")[-4:]} for a in accounts_raw]
        txs = []
        for a in accounts_raw:
            iban = a.get("iban")
            if not iban:
                continue
            try:
                d = await asyncio.to_thread(_qonto_get, "transactions", {"iban": iban, "per_page": 100})
                for t in (d.get("transactions", []) or []):
                    txs.append({"account": a.get("name") or a.get("slug"),
                                "date": t.get("settled_at") or t.get("emitted_at"),
                                "amount": float(t.get("amount") or 0), "side": t.get("side"),
                                "label": t.get("label") or t.get("clean_counterparty_name") or t.get("note") or "",
                                "type": t.get("operation_type")})
            except Exception as e:
                logger.warning(f"neo: Qonto transactions ({(iban or '')[-4:]}) : {e}")
        txs.sort(key=lambda x: x.get("date") or "", reverse=True)
        # Agrégats : mensuel (revenus/dépenses) + répartition des dépenses par type
        from collections import defaultdict
        monthly = {}
        by_type = defaultdict(float)
        for t in txs:
            m = (t.get("date") or "")[:7]
            if not m:
                continue
            # On garde TOUS les types (virements inclus) : un virement reçu = vraie entrée (client),
            # un virement émis = vraie sortie (fournisseur/salaire). Seuls les mouvements internes
            # entre comptes propres gonfleraient les totaux, mais les sous-comptes sont inactifs.
            slot = monthly.setdefault(m, {"month": m, "income": 0.0, "expense": 0.0})
            amt = t.get("amount") or 0
            if t.get("side") == "credit":
                slot["income"] += amt
            else:
                slot["expense"] += amt
                by_type[t.get("type") or "autre"] += amt
        months_sorted = sorted(monthly.values(), key=lambda x: x["month"])[-6:]
        for mm in months_sorted:
            mm["income"], mm["expense"] = round(mm["income"], 2), round(mm["expense"], 2)
        by_type_list = sorted(({"type": k, "amount": round(v, 2)} for k, v in by_type.items()),
                              key=lambda x: x["amount"], reverse=True)
        cur = monthly.get(datetime.now(timezone.utc).strftime("%Y-%m"), {"income": 0, "expense": 0})
        return {"success": True, "connected": True, "total_balance": round(sum(x["balance"] for x in accounts), 2),
                "currency": "EUR", "accounts": accounts, "transactions": txs[:limit],
                "monthly": months_sorted, "by_type": by_type_list, "tx_count": len(txs),
                "month_income": round(cur.get("income", 0), 2), "month_expense": round(cur.get("expense", 0), 2)}
    except Exception as e:
        logger.warning(f"neo: Qonto trésorerie indisponible: {e}")
        return {"success": True, "connected": False, "error": str(e)[:200]}


# ==================== Health score (Phase 4) ====================
async def _compute_health_score() -> dict:
    """Score de santé business /100 : trésorerie + impayés + pipeline (devis) + leads chauds non traités."""
    now = datetime.now(timezone.utc)
    month = now.strftime("%Y-%m")
    try:
        entries = await db.budget.find({"date": {"$regex": f"^{month}"}}, {"_id": 0, "amount": 1, "type": 1}).to_list(500)
    except Exception:
        entries = []
    income = sum((e.get("amount", 0) or 0) for e in entries if e.get("type") == "income")
    expense = sum((e.get("amount", 0) or 0) for e in entries if e.get("type") == "expense")
    balance = income - expense
    invoices = await db.invoices.find({}, {"_id": 0, "document_type": 1, "status": 1, "total": 1, "due_date": 1}).to_list(2000)
    PAID = ("payée", "payee", "payé", "paye", "annulée", "annulee", "annulé", "brouillon")
    overdue = [i for i in invoices if i.get("document_type") != "devis"
               and (i.get("status") or "").lower() not in PAID
               and (i.get("status") == "en_retard" or (_to_dt(i.get("due_date")) and _to_dt(i.get("due_date")) < now))]
    pending_devis = [i for i in invoices if i.get("document_type") == "devis" and (i.get("status") or "").lower() == "brouillon"]
    contacts = await db.contacts.find({"status": "nouveau"}, {"_id": 0, "score": 1, "score_value": 1}).to_list(2000)

    def _hot(c):
        sv = c.get("score_value")
        return (sv >= 70) if isinstance(sv, (int, float)) else (c.get("score") in ("chaud", "chaude"))
    hot = [c for c in contacts if _hot(c)]
    score = 70
    score += 15 if balance >= 0 else -20
    score -= min(35, len(overdue) * 12)
    score += min(12, len(pending_devis) * 3)
    score -= min(15, len(hot) * 5)
    score = max(0, min(100, int(round(score))))
    label = "solide" if score >= 75 else ("à surveiller" if score >= 50 else "tendu")
    return {"success": True, "score": score, "label": label, "balance": round(balance, 2),
            "overdue_count": len(overdue), "overdue_total": round(sum((i.get("total") or 0) for i in overdue), 2),
            "pending_devis": len(pending_devis), "hot_leads": len(hot)}


# ==================== Budget incarné (Phase 2 — VOIR l'argent) ====================
async def _budget_context() -> str:
    """État financier EN DIRECT, injecté dans le cerveau de Néo à chaque message : il « voit »
    la trésorerie sans avoir à appeler un outil. Données locales Mongo uniquement (pas d'appel
    réseau Qonto ici — gardé pour l'outil get_bank_balance à la demande). Robuste : jamais bloquant."""
    try:
        now = datetime.now(timezone.utc)
        month = now.strftime("%Y-%m")
        entries = await db.budget.find({"date": {"$regex": f"^{month}"}},
                                       {"_id": 0, "amount": 1, "type": 1}).to_list(500)
        income = sum((e.get("amount", 0) or 0) for e in entries if e.get("type") == "income")
        expense = sum((e.get("amount", 0) or 0) for e in entries if e.get("type") == "expense")
        balance = income - expense
        ytd = await db.budget.find({"date": {"$regex": f"^{now.strftime('%Y')}"}, "type": "income"},
                                   {"_id": 0, "amount": 1}).to_list(3000)
        income_ytd = sum((e.get("amount", 0) or 0) for e in ytd)
        invoices = await db.invoices.find({}, {"_id": 0, "document_type": 1, "status": 1,
                                               "total": 1, "due_date": 1}).to_list(2000)
        PAID = ("payée", "payee", "payé", "paye", "annulée", "annulee", "annulé", "brouillon")
        overdue = [i for i in invoices if i.get("document_type") != "devis"
                   and (i.get("status") or "").lower() not in PAID
                   and (i.get("status") == "en_retard" or (_to_dt(i.get("due_date")) and _to_dt(i.get("due_date")) < now))]
        overdue_total = sum((i.get("total") or 0) for i in overdue)
        pending_devis = [i for i in invoices if i.get("document_type") == "devis"
                         and (i.get("status") or "").lower() in ("brouillon", "en_attente", "envoyée", "envoyee")]
        pipeline_total = sum((i.get("total") or 0) for i in pending_devis)
        try:
            fc = await db.budget_forecasts.find({"month": month},
                                                {"_id": 0, "planned_amount": 1, "type": 1}).to_list(500)
        except Exception:
            fc = []
        planned_income = sum((f.get("planned_amount", 0) or 0) for f in fc if f.get("type") == "income")
        planned_expense = sum((f.get("planned_amount", 0) or 0) for f in fc if f.get("type") == "expense")
    except Exception as e:
        logger.warning(f"neo budget_context KO: {e}")
        return ""
    lines = [f"\n\n=== ÉTAT FINANCIER EN DIRECT ({month}) — tiens-en compte sans le redemander, cite ces chiffres quand c'est utile ==="]
    lines.append(f"Réel ce mois : encaissé {income:.0f}€, dépensé {expense:.0f}€, solde {balance:+.0f}€.")
    lines.append(f"Encaissé cumulé depuis janvier {month[:4]} : {income_ytd:.0f}€.")
    if planned_income or planned_expense:
        lines.append(f"Prévisionnel saisi (page Budget) : {planned_income:.0f}€ de recettes prévues, {planned_expense:.0f}€ de dépenses prévues.")
    lines.append(f"Impayés à relancer : {len(overdue)} facture(s) pour {overdue_total:.0f}€." if overdue
                 else "Impayés à relancer : aucun.")
    if pending_devis:
        lines.append(f"Pipeline (devis en attente) : {len(pending_devis)} pour {pipeline_total:.0f}€ potentiels.")
        lines.append(f"Atterrissage du mois si ces devis passent : {income + pipeline_total:.0f}€ encaissés.")
    return "\n".join(lines)


# ==================== Pilotage stratégique (Phase 5) ====================
async def _strategic_review(args=None, uid=None) -> dict:
    """Point stratégique : prévisionnel (encaissé + pipeline), risque n°1, reco 'prêt à déléguer'.
    Jugement routé vers Claude (repli Gemini). Lecture seule."""
    hs = await _compute_health_score()
    invoices = await db.invoices.find({}, {"_id": 0, "document_type": 1, "status": 1, "total": 1}).to_list(2000)
    pending_devis = [i for i in invoices if i.get("document_type") == "devis"
                     and (i.get("status") or "").lower() in ("brouillon", "en_attente", "envoyée", "envoyee")]
    pipeline_total = sum((i.get("total") or 0) for i in pending_devis)
    try:
        open_tasks = await db.tasks.count_documents({"status": {"$nin": ["done", "cancelled"]}})
        total_contacts = await db.contacts.count_documents({})
        clients = await db.contacts.count_documents({"status": "client"})
        objs = [r.get("content") for r in await db.neo_memory.find(
            {"type": "objective"}, {"_id": 0, "content": 1}).to_list(10)]
    except Exception:
        open_tasks = total_contacts = clients = 0
        objs = []
    data = (f"Trésorerie ce mois (solde): {hs['balance']}€. Impayés à relancer: {hs['overdue_count']} "
            f"({hs['overdue_total']}€). Pipeline (devis en attente): {len(pending_devis)} pour {pipeline_total:.0f}€ "
            f"potentiels. Leads chauds non traités: {hs['hot_leads']}. Tâches ouvertes: {open_tasks}. "
            f"Contacts: {total_contacts}, dont {clients} clients. Score santé: {hs['score']}/100. "
            f"Objectifs enregistrés: {('; '.join(o for o in objs if o)) if objs else 'aucun (à définir avec Léo)'}.")
    system = ("Tu es Néo, l'associé co-gérant IA d'Alpha Agency (agence de communication, Guadeloupe). Ta raison "
              "d'être : faire croître le CA, la marge et le bénéfice. Donne un POINT STRATÉGIQUE concis et "
              "actionnable, en français, sans tirets longs : 1) prévisionnel du mois (encaissé + pipeline, vs "
              "objectif si connu, alerte si dérapage) ; 2) le risque n°1 ou là où l'on perd des prospects ; "
              "3) reco « prêt à déléguer ? » en comparant la charge (tâches) et la trésorerie. 4 à 7 phrases, direct, "
              "termine par la prochaine action concrète à faire.")
    review, model = await _strategic_text(system, "Données actuelles du CRM :\n" + data + "\n\nFais ton point stratégique.")
    return {"success": True, "model": model, "review": review, "pipeline_total": round(pipeline_total, 2),
            "open_tasks": open_tasks, "score": hs["score"], "balance": hs["balance"]}


# ==================== Registre d'outils (Partie C) ====================
def _obj(props: dict, required=None):
    return {"type": "object", "properties": props, "required": required or []}


_STR = {"type": "string"}
_INT = {"type": "integer"}
_BOOL = {"type": "boolean"}
_NUM = {"type": "number"}
# Ligne de devis/facture, avec remise EXPLICITE par ligne (le moteur sait la calculer + l'afficher).
_QUOTE_ITEM = {"type": "object", "properties": {
    "title": _STR, "description": _STR, "quantity": _NUM, "unit_price": _NUM,
    "discount": _NUM, "discountType": {"type": "string", "enum": ["%", "€"]},
}}

# Chaque outil : name, description, params (json schema), validation (garde-fou A.5), run(args, uid)
async def _exec_web_search(args, uid):
    """Recherche web temps réel via Gemini (Google Search grounding), appel séparé du loop d'outils."""
    q = (args.get("query") or "").strip()
    if not q:
        return {"success": False, "message": "Requête vide."}
    if not _client:
        return {"success": False, "message": "Recherche indisponible (IA non configurée)."}
    def _call():
        cfg = _t.GenerateContentConfig(tools=[_t.Tool(google_search=_t.GoogleSearch())])
        return _client.models.generate_content(model="gemini-2.5-flash", contents=q, config=cfg)
    try:
        resp = await asyncio.to_thread(_call)
        txt = (getattr(resp, "text", "") or "").strip()
        sources = []
        try:
            gm = resp.candidates[0].grounding_metadata
            for c in (getattr(gm, "grounding_chunks", None) or [])[:5]:
                w = getattr(c, "web", None)
                if w:
                    sources.append({"title": getattr(w, "title", "") or "", "uri": getattr(w, "uri", "") or ""})
        except Exception:
            pass
        return {"success": True, "result": txt[:4000] or "Aucun résultat.", "sources": sources}
    except Exception as e:
        logger.warning(f"neo web_search: {e}")
        return {"success": False, "message": f"Recherche échouée: {str(e)[:150]}"}


async def _exec_get_invoice_pdf(args, uid):
    """Récupère un devis/facture (par numéro ou client) et renvoie un lien de téléchargement (30 min)."""
    import secrets
    num = (args.get("invoice_number") or "").strip()
    name = (args.get("contact_name") or "").strip()
    inv = None
    if num:
        inv = await db.invoices.find_one({"invoice_number": {"$regex": num, "$options": "i"}}, {"_id": 0})
    if not inv and name:
        c = await _find_contact(name)
        if c:
            rows = await db.invoices.find({"contact_id": c.get("id")}, {"_id": 0}).sort("created_at", -1).to_list(1)
            inv = rows[0] if rows else None
    if not inv:
        return {"success": False, "message": "Aucun devis/facture trouvé pour cette référence ou ce client."}
    token = secrets.token_urlsafe(32)
    await db.pdf_tokens.insert_one({
        "token": token, "invoice_id": inv["id"],
        "expires_at": (datetime.now(timezone.utc) + timedelta(minutes=30)).isoformat(),
        "created_at": datetime.now(timezone.utc).isoformat(),
    })
    dtype = "devis" if inv.get("document_type") == "devis" else "facture"
    url = f"{PUBLIC_BASE}/api/invoices/{inv['id']}/pdf-download/{token}"
    return {"success": True,
            "result": {"type": dtype, "number": inv.get("invoice_number"), "total": inv.get("total"), "download_url": url},
            "message": f"{dtype} {inv.get('invoice_number')} prêt à télécharger : {url}"}


async def _exec_delegate_task(args, uid):
    """Délègue une sous-tâche lourde à un sous-agent focalisé (rédaction, analyse, idéation, plan)."""
    objective = (args.get("objective") or "").strip()
    context = (args.get("context") or "").strip()
    if not objective:
        return {"success": False, "message": "Objectif vide."}
    sub_system = ("Tu es un sous-agent expert d'Alpha Agency (agence de communication digitale, Guadeloupe), "
                  "missionné par Néo. Réalise la tâche de façon complète, concrète et directement exploitable. "
                  "Sois structuré et actionnable. Réponds en français.")
    prompt = f"Mission : {objective}"
    if context:
        prompt += f"\n\nContexte :\n{context}"
    try:
        out = await _gemini_text(sub_system, prompt)
        return {"success": True, "result": (out or "")[:6000]}
    except Exception:
        try:
            out = await _claude_text(sub_system, prompt)
            return {"success": True, "result": (out or "")[:6000]}
        except Exception as e2:
            return {"success": False, "message": f"Sous-agent indisponible: {str(e2)[:120]}"}


async def _exec_send_to_cowork(args, uid):
    """Néo envoie une tâche/brief à Claude Cowork (sur le PC de Léo) via la boîte cowork_inbox.
    Cowork la récupère ensuite (outil MCP get_cowork_tasks) et la traite."""
    title = (args.get("title") or "").strip()
    brief = (args.get("brief") or "").strip()
    if not (title or brief):
        return {"success": False, "message": "Titre ou brief requis."}
    doc = {"id": str(uuid.uuid4()), "title": title or "Tâche", "brief": brief, "status": "pending",
           "created_at": datetime.now(timezone.utc).isoformat(), "user_id": uid}
    await db.cowork_inbox.insert_one(doc)
    return {"success": True, "result": {"id": doc["id"], "title": doc["title"]},
            "message": f"Tâche déposée pour Cowork : « {doc['title']} ». Le Claude de Léo sur PC la récupérera."}


async def _exec_generate_image(args, uid):
    """Génère une image/visuel via Gemini et renvoie une URL affichable (stockée en base)."""
    prompt = (args.get("prompt") or "").strip()
    if not prompt:
        return {"success": False, "message": "Décris l'image à générer."}
    if not _client:
        return {"success": False, "message": "Génération d'image indisponible (IA non configurée)."}
    import base64 as _b64
    # Chaîne de modèles image (le 1er accessible gagne) — robuste aux renommages/dispos de la clé
    models = list(dict.fromkeys([m for m in [
        NEO_IMAGE_MODEL, "gemini-2.5-flash-image-preview", "gemini-2.5-flash-image",
        "gemini-2.0-flash-exp-image-generation", "imagen-3.0-generate-002",
    ] if m]))
    raw, mime, last_err = None, "image/png", "?"
    for mdl in models:
        def _call(_m=mdl):
            cfg = _t.GenerateContentConfig(response_modalities=["TEXT", "IMAGE"])
            return _client.models.generate_content(model=_m, contents=prompt, config=cfg)
        try:
            resp = await asyncio.to_thread(_call)
            for part in (resp.candidates[0].content.parts or []):
                inl = getattr(part, "inline_data", None)
                if inl and getattr(inl, "data", None):
                    d = inl.data
                    raw = d if isinstance(d, (bytes, bytearray)) else _b64.b64decode(d)
                    mime = getattr(inl, "mime_type", "image/png") or "image/png"
                    break
            if raw:
                break
            last_err = "réponse sans image"
        except Exception as e:
            last_err = str(e)[:120]
            continue
    if not raw:
        logger.warning(f"neo generate_image: aucun modèle ({last_err})")
        return {"success": False, "message": f"Génération d'image indisponible (modèle image absent de la clé : {last_err})."}
    img_id = str(uuid.uuid4())
    await db.neo_images.insert_one({"id": img_id, "data": _b64.b64encode(raw).decode(), "mime": mime,
                                    "prompt": prompt[:500], "user_id": uid,
                                    "created_at": datetime.now(timezone.utc).isoformat()})
    url = f"{PUBLIC_BASE}/api/neo/image/{img_id}"
    return {"success": True, "result": {"image_url": url}, "message": f"Visuel généré : {url}"}


async def _exec_create_contact(args, uid):
    """Crée une fiche contact dans db.contacts (même schéma que la route POST /contacts)."""
    first = (args.get("first_name") or "").strip()
    last = (args.get("last_name") or "").strip()
    if not first and not last:
        name = (args.get("name") or "").strip()
        if name:
            parts = name.split()
            first = parts[0]
            last = " ".join(parts[1:])
    if not first and not last:
        return {"success": False, "message": "Donne au moins un nom pour créer la fiche."}
    contact_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    doc = {"id": contact_id, "first_name": first, "last_name": last,
           "email": (args.get("email") or "").strip() or None,
           "phone": (args.get("phone") or "").strip() or None,
           "company": (args.get("company") or "").strip() or None,
           "position": (args.get("position") or "").strip() or None,
           "city": (args.get("city") or "").strip() or None,
           "notes": (args.get("notes") or "").strip() or None,
           "source": (args.get("source") or "Néo").strip(),
           "status": "nouveau", "score": "tiède",
           "created_at": now, "updated_at": now, "created_by": uid}
    await db.contacts.insert_one(doc)
    full = f"{first} {last}".strip()
    return {"success": True, "result": {"contact_id": contact_id, "name": full},
            "message": f"Fiche contact créée : {full}."}


async def _exec_create_quote(args, uid):
    """Crée un VRAI devis dans db.invoices (document_type=devis) -> visible dans la page Facturation
    ET compté dans le pipeline/prévisionnel. Lié au contact si on le retrouve. Reste un brouillon (pas d'envoi)."""
    from .invoices import get_next_invoice_number, calculate_invoice_totals  # lazy (évite la circularité)
    contact = None
    if args.get("contact_id") or args.get("contact_name") or args.get("client_name"):
        try:
            contact = await _find_contact({"contact_id": args.get("contact_id"),
                                           "contact_name": args.get("contact_name") or args.get("client_name")})
        except Exception:
            contact = None
    contact_id = contact.get("id") if contact else None
    raw = args.get("services") or args.get("items") or []

    def _num(v, d=0.0):
        try:
            return float(v)
        except (TypeError, ValueError):
            return d
    items = []
    for s in (raw if isinstance(raw, list) else []):
        if not isinstance(s, dict):
            continue
        up = s.get("unit_price")
        if up is None:
            up = s.get("price")
        dtype = (s.get("discountType") or "%")
        dtype = "%" if dtype in ("%", "percent") else "€"
        items.append({"title": s.get("title") or s.get("name") or "",
                      "description": s.get("description") or s.get("title") or "Prestation",
                      "quantity": _num(s.get("quantity"), 1) or 1,
                      "unit_price": _num(up),
                      "discount": _num(s.get("discount")),       # remise par ligne (jamais bakée dans le prix)
                      "discountType": dtype})
    if not items:
        return {"success": False, "message": "Décris au moins une ligne de prestation (services) pour le devis."}
    gd = _num(args.get("global_discount"))
    gdt = (args.get("global_discount_type") or "%")
    gdt = "%" if gdt in ("%", "percent") else "€"
    subtotal, tva, total = calculate_invoice_totals(items, gd, gdt)
    number = await get_next_invoice_number("devis", "standard", None)
    now = datetime.now(timezone.utc)
    client_name = (args.get("client_name")
                   or (f"{contact.get('first_name','')} {contact.get('last_name','')}".strip() if contact else "")) or "Client"
    doc = {"id": str(uuid.uuid4()), "invoice_number": number, "quote_id": None, "contact_id": contact_id,
           "document_type": "devis", "invoice_type": "standard", "parent_invoice_id": None, "parent_invoice_number": None,
           "items": items, "subtotal": subtotal, "tva": tva, "total": total, "total_paid": 0, "remaining": total,
           "globalDiscount": gd, "globalDiscountType": gdt, "status": "brouillon",
           "due_date": (now + timedelta(days=30)).strftime("%Y-%m-%d"), "payment_terms": "30",
           "notes": args.get("notes"), "conditions": None, "client_name": client_name,
           "created_at": now.isoformat(), "created_by": uid}
    await db.invoices.insert_one(doc)
    if contact:
        extra = contact.get("email") or contact.get("company") or ""
        link = f"rattaché au contact « {client_name} »" + (f" ({extra})" if extra else "")
    else:
        link = "⚠️ AUCUN contact existant rattaché — vérifie le nom ou crée d'abord la fiche"
    has_line_remise = any(it["discount"] for it in items)
    if gd:
        remise_txt = f" Remise globale de {gd:.0f}{'%' if gdt == '%' else '€'} appliquée (le devis affiche le prix d'origine + la remise)."
    elif has_line_remise:
        remise_txt = " Remise(s) par ligne appliquée(s) (le devis affiche le prix d'origine + la remise)."
    else:
        remise_txt = ""
    return {"success": True, "result": {"quote_id": doc["id"], "number": number, "total": total,
                                        "subtotal_ht": round(subtotal, 2), "global_discount": gd,
                                        "contact_id": contact_id, "client_name": client_name},
            "message": f"Devis {number} créé — total TTC {total:.2f}€ (HT après remise {subtotal:.2f}€).{remise_txt} {link}. "
                       f"Brouillon visible dans Facturation. Indique à Léo le contact EXACT rattaché et demande si c'est le bon."}


TOOLS = [
    # --- Lecture ---
    {"name": "web_search", "validation": False, "run": _exec_web_search,
     "description": "Recherche sur le web en temps réel (Google via Gemini) : prix du marché, info sur une entreprise/personne, veille concurrentielle, tendances, actualité, inspiration créative. À utiliser dès qu'une info récente ou externe au CRM est utile.",
     "params": _obj({"query": _STR})},
    {"name": "get_invoice_pdf", "validation": False, "run": _exec_get_invoice_pdf,
     "description": "Récupère un devis ou une facture (par numéro ex 'FAC-2026-0010', ou par nom de client = le plus récent) et renvoie un LIEN de téléchargement du PDF à donner à Léo. Utilise-le quand Léo demande de lui envoyer/récupérer un devis ou une facture.",
     "params": _obj({"invoice_number": _STR, "contact_name": _STR})},
    {"name": "delegate_task", "validation": False, "run": _exec_delegate_task,
     "description": "Délègue une sous-tâche lourde à un sous-agent focalisé : rédiger une proposition/offre commerciale, analyser un marché, générer des idées de contenu/campagne, structurer un plan. Donne un objectif clair + le contexte utile. À utiliser quand la tâche mérite un vrai travail dédié.",
     "params": _obj({"objective": _STR, "context": _STR})},
    {"name": "send_to_cowork", "validation": False, "run": _exec_send_to_cowork,
     "description": "Envoie une tâche / un brief à Claude Cowork sur le PC de Léo (pour qu'il bosse ou réfléchisse dessus côté code/dev). Donne un titre court + un brief détaillé. Utilise-le quand Léo dit d'envoyer/confier quelque chose à Cowork ou à son Claude sur PC.",
     "params": _obj({"title": _STR, "brief": _STR})},
    {"name": "generate_image", "validation": False, "run": _exec_generate_image,
     "description": "Génère une IMAGE / un visuel à partir d'une description détaillée (visuel de post réseaux, illustration, concept créatif, mockup). Renvoie une URL d'image à montrer à Léo. Utilise-le dès que Léo demande de créer/générer une image, un visuel ou une illustration.",
     "params": _obj({"prompt": _STR})},
    {"name": "search_contacts", "validation": False, "run": _exec_search_contacts,
     "description": "Cherche des contacts par texte (nom/entreprise/email) et/ou statut.",
     "params": _obj({"query": _STR, "status": _STR, "limit": _INT})},
    {"name": "get_contact", "validation": False, "run": _exec_get_contact,
     "description": "Récupère la fiche complète d'un contact par nom/entreprise/email ou id.",
     "params": _obj({"contact_name": _STR, "contact_id": _STR})},
    {"name": "get_contact_history", "validation": False, "run": _exec_get_contact_history,
     "description": "Historique chronologique d'un contact (lead, relances, devis/factures, tâches) — pour raconter son parcours.",
     "params": _obj({"contact_name": _STR, "contact_id": _STR})},
    {"name": "read_emails", "validation": False, "run": _exec_read_emails,
     "description": "Lit la boîte Gmail. 'query' = recherche Gmail (ex: 'from:client@x.com', 'is:unread', 'newer_than:7d', un nom). Pour croiser les mails avec les contacts/impayés ou préparer un échange.",
     "params": _obj({"query": _STR, "limit": _INT})},
    {"name": "list_leads", "validation": False, "run": _exec_list_leads,
     "description": "Liste les leads récents (site + chatbot). only_hot=true pour les leads chauds (score>=70).",
     "params": _obj({"only_hot": _BOOL, "limit": _INT})},
    {"name": "get_budget_summary", "validation": False, "run": _exec_get_budget_summary,
     "description": "Résumé financier du mois : entrées, sorties, solde, impayés (nb + total), devis en attente, + solde bancaire réel Qonto si connecté.",
     "params": _obj({})},
    {"name": "get_bank_balance", "validation": False, "run": lambda a, u: _qonto_summary(),
     "description": "Solde bancaire RÉEL via Qonto (total + comptes). Pour répondre à la trésorerie réelle de l'agence.",
     "params": _obj({})},
    {"name": "list_transactions", "validation": False, "run": lambda a, u: _qonto_treasury(int(a.get("limit") or 25)),
     "description": "Transactions bancaires récentes (Qonto) : virements, cartes, prélèvements, encaissements + soldes par compte. Pour analyser les mouvements de trésorerie.",
     "params": _obj({"limit": _INT})},
    {"name": "list_overdue_invoices", "validation": False, "run": _exec_list_overdue_invoices,
     "description": "Liste les factures en retard (à relancer) avec le total dû.",
     "params": _obj({})},
    {"name": "list_tasks", "validation": False, "run": _exec_list_tasks,
     "description": "Liste les tâches (par défaut les non terminées). status optionnel.",
     "params": _obj({"status": _STR})},
    {"name": "get_health_score", "validation": False, "run": lambda a, u: _compute_health_score(),
     "description": "Score de santé business /100 (trésorerie, impayés, devis en attente, leads chauds) + détail.",
     "params": _obj({})},
    {"name": "strategic_review", "validation": False, "run": lambda a, u: _strategic_review(a, u),
     "description": "Point stratégique (jugement routé vers Claude) : prévisionnel du mois, risque n°1, reco 'prêt à déléguer'. À utiliser quand Léo demande un point stratégique / prévisionnel / conseil de pilotage.",
     "params": _obj({})},
    {"name": "activity_report", "validation": False, "run": lambda a, u: activity_report_action(a),
     "description": "Reporting d'activité : leads, conversion, valeur moyenne des devis, service & canal n°1.",
     "params": _obj({"days": _INT})},
    {"name": "prep_meeting", "validation": False, "run": lambda a, u: prep_meeting_action(a),
     "description": "Prépare un RDV : fiche + devis/factures + tâches + points à valider.",
     "params": _obj({"contact_name": _STR, "contact_id": _STR})},
    {"name": "enrich_company", "validation": False, "run": lambda a, u: enrich_company_action(a),
     "description": "Recherche web sur l'entreprise d'un contact (avant un appel) + 3 angles de RDV. Stocke un résumé.",
     "params": _obj({"contact_name": _STR, "contact_id": _STR})},
    {"name": "list_documents", "validation": False, "run": lambda a, u: list_documents_action(a),
     "description": "Liste les documents (filtres optionnels).",
     "params": _obj({"folder_name": _STR, "file_type": _STR, "search": _STR})},
    {"name": "get_document", "validation": False, "run": lambda a, u: get_document_action(a),
     "description": "Détails d'un document par nom ou id.",
     "params": _obj({"document_name": _STR, "document_id": _STR})},
    {"name": "crm_query", "validation": False, "run": _exec_crm_query,
     "description": "Lecture seule générique sur une collection du CRM pour répondre à une question imprévue. "
                    "collections: contacts, invoices, tasks, appointments, budget, opportunities, multilink_pages, documents, quotes, portfolio.",
     "params": _obj({"collection": _STR, "filter": {"type": "object"}, "limit": _INT}, ["collection"])},
    # --- Écriture générique encadrée : couverture TOTALE du CRM (tout ce qui n'a pas d'outil dédié) ---
    {"name": "crm_create", "validation": False, "run": _exec_crm_create,
     "description": "Crée un document dans N'IMPORTE QUELLE collection métier du CRM (couverture totale). Utilise-le pour toute création sans outil dédié : opportunité, RDV, service, dossier, document, tâche, colonne de pipeline, séquence de nurturing, post éditorial, etc. Args: collection + data (objet des champs). Collections permises: contacts, invoices, quotes, tasks, opportunities, appointments, documents, folders, tags, portfolio, services, notes, pipeline_columns, budget, budget_forecasts, budget_categories, multilink_*, editorial_*, blog_posts, news_articles, scheduled_posts, nurturing_*.",
     "params": _obj({"collection": _STR, "data": {"type": "object"}}, ["collection", "data"])},
    {"name": "crm_update", "validation": False, "run": _exec_crm_update,
     "description": "Modifie un document du CRM (couverture totale). Args: collection + filter (cible, ex {\"id\":\"...\"}) + updates (champs à changer). Pour toute modification sans outil dédié (changer un statut, un montant, une date, des champs...).",
     "params": _obj({"collection": _STR, "filter": {"type": "object"}, "updates": {"type": "object"}}, ["collection", "filter", "updates"])},
    {"name": "crm_delete", "validation": True, "run": _exec_crm_delete,
     "description": "Supprime un document du CRM. Args: collection + filter ciblant le document. IRRÉVERSIBLE -> validation de Léo requise.",
     "params": _obj({"collection": _STR, "filter": {"type": "object"}}, ["collection", "filter"])},
    # --- Actions internes (pas de sortie client → exécution directe, journalisée) ---
    {"name": "create_task", "validation": False, "run": lambda a, u: create_task_action(a, u),
     "description": "Crée une tâche. title requis ; due_date YYYY-MM-DD ; priority low/medium/high/urgent.",
     "params": _obj({"title": _STR, "description": _STR, "priority": _STR, "due_date": _STR, "contact_id": _STR}, ["title"])},
    {"name": "mark_task_done", "validation": False, "run": lambda a, u: mark_task_done_action(a),
     "description": "Marque une tâche comme terminée (task_title ou task_id).",
     "params": _obj({"task_title": _STR, "task_id": _STR})},
    {"name": "schedule_followup", "validation": False, "run": lambda a, u: schedule_followup_action(a, u),
     "description": "Programme une relance (crée une tâche catégorie relance).",
     "params": _obj({"contact_name": _STR, "contact_id": _STR, "due_date": _STR, "label": _STR})},
    {"name": "create_contact", "validation": False, "run": _exec_create_contact,
     "description": "Crée une fiche contact dans le CRM (nouveau prospect/client). 'first_name' requis (ou 'name' = nom complet) ; optionnels : last_name, email, phone, company, position, city, notes. Quand Léo veut créer un client puis un devis, crée D'ABORD le contact ici, puis enchaîne create_quote.",
     "params": _obj({"first_name": _STR, "last_name": _STR, "name": _STR, "email": _STR, "phone": _STR, "company": _STR, "position": _STR, "city": _STR, "notes": _STR})},
    {"name": "set_contact_status", "validation": False, "run": lambda a, u: set_contact_status_action(a),
     "description": "Change le statut d'un contact (gagné/perdu/en cours/client/qualifié...).",
     "params": _obj({"contact_name": _STR, "contact_id": _STR, "status": _STR}, ["status"])},
    {"name": "add_contact_note", "validation": False, "run": lambda a, u: add_contact_note_action(a),
     "description": "Ajoute une note horodatée à une fiche (sans rien écraser).",
     "params": _obj({"contact_name": _STR, "contact_id": _STR, "note": _STR}, ["note"])},
    {"name": "update_contact", "validation": False, "run": lambda a, u: update_contact_action(a),
     "description": "Met à jour des champs d'un contact (phone, email, company, budget, poste, project_type, city, note...).",
     "params": _obj({"contact_name": _STR, "contact_id": _STR, "updates": {"type": "object"}}, ["updates"])},
    {"name": "create_quote", "validation": False, "run": _exec_create_quote,
     "description": "Crée un VRAI devis (brouillon) dans la FACTURATION : visible dans la page Facturation et compté dans le pipeline/prévisionnel. Lie-le à un contact existant via contact_name (si le client n'existe pas, crée d'abord la fiche avec create_contact). "
                    "services: [{title, description, quantity, unit_price, discount, discountType}]. "
                    "REMISE — TRÈS IMPORTANT : pour une réduction, mets le PRIX NORMAL (avant remise) dans unit_price et le montant de la remise dans discount, avec discountType '€' (montant en euros) ou '%' (pourcentage). "
                    "Ne calcule JAMAIS le prix remisé toi-même pour le mettre dans unit_price, et n'explique PAS la remise dans la description : le devis affiche déjà une colonne remise et recalcule le total. "
                    "Pour une remise sur le TOTAL du devis (et non une ligne), utilise global_discount + global_discount_type. "
                    "Ne l'envoie pas (l'envoi est une étape séparée avec validation).",
     "params": _obj({"contact_name": _STR, "client_name": _STR, "client_email": _STR,
                     "services": {"type": "array", "items": _QUOTE_ITEM},
                     "global_discount": _NUM, "global_discount_type": {"type": "string", "enum": ["%", "€"]},
                     "notes": _STR})},
    {"name": "draft_followup_email", "validation": False, "run": lambda a, u: draft_followup_email_action(a),
     "description": "Rédige (sans envoyer) un email de relance personnalisé, stocké en brouillon.",
     "params": _obj({"contact_name": _STR, "contact_id": _STR, "angle": _STR})},
    # --- Mémoire de Néo (Phase 3) ---
    {"name": "remember", "validation": False, "run": _exec_remember,
     "description": "Mémorise un fait durable. type: objective|rule|daily_log|client_fact|decision ; content requis.",
     "params": _obj({"type": _STR, "content": _STR, "key": _STR, "contact_id": _STR}, ["content"])},
    {"name": "recall", "validation": False, "run": _exec_recall,
     "description": "Retrouve dans la mémoire de Néo (filtre par texte 'query' et/ou 'type').",
     "params": _obj({"query": _STR, "type": _STR})},
    {"name": "update_objective", "validation": False, "run": _exec_update_objective,
     "description": "Définit/met à jour un objectif chiffré de l'agence (key optionnel, défaut objectif_principal).",
     "params": _obj({"content": _STR, "key": _STR}, ["content"])},
    {"name": "log_day", "validation": False, "run": _exec_log_day,
     "description": "Ajoute une entrée au journal du jour (check-in matin/soir, avancement de Léo).",
     "params": _obj({"note": _STR}, ["note"])},
    # --- Actions SORTANTES / IRRÉVERSIBLES → validation humaine (garde-fou A.5) ---
    {"name": "send_followup", "validation": True, "run": lambda a, u: send_followup_action(a),
     "description": "ENVOIE la relance (email) au prospect. Sortie client → nécessite la validation de Léo.",
     "params": _obj({"contact_name": _STR, "contact_id": _STR})},
    {"name": "merge_contacts", "validation": True, "run": lambda a, u: merge_contacts_action(a),
     "description": "Fusionne deux fiches en doublon (supprime le doublon). Irréversible → validation de Léo.",
     "params": _obj({"keep": _STR, "remove": _STR, "primary_id": _STR, "duplicate_id": _STR})},
]
_SPEC = {t["name"]: t for t in TOOLS}


def _gemini_tools():
    if not _t:
        return None
    decls = [_t.FunctionDeclaration(name=t["name"], description=t["description"],
                                    parameters_json_schema=t["params"]) for t in TOOLS]
    return [_t.Tool(function_declarations=decls)]


# ==================== Journalisation (Phase 0) ====================
async def _log(kind: str, payload: dict):
    try:
        await db.neo_action_log.insert_one({
            "id": str(uuid.uuid4()), "kind": kind,
            "at": datetime.now(timezone.utc).isoformat(), **payload})
    except Exception as e:
        logger.warning(f"neo log échoué: {e}")


# ==================== Exécution d'un outil (avec garde-fous) ====================
async def execute_tool(name: str, args: dict, user_id: str, confirmed: bool = False) -> dict:
    spec = _SPEC.get(name)
    if not spec:
        return {"success": False, "error": f"Outil inconnu: {name}"}
    # Garde-fou A.5 : sortie client / suppression → validation humaine
    if spec["validation"] and not confirmed:
        action_id = str(uuid.uuid4())
        await db.neo_pending_actions.insert_one({
            "id": action_id, "name": name, "args": args, "user_id": user_id,
            "created_at": datetime.now(timezone.utc).isoformat(), "status": "pending"})
        await _log("pending", {"tool": name, "args": args, "user_id": user_id, "action_id": action_id})
        return {"success": True, "pending": True, "action_id": action_id,
                "message": f"EN ATTENTE DE VALIDATION de Léo (action: {name})."}
    try:
        res = await spec["run"](args, user_id)
    except Exception as e:
        logger.error(f"neo tool {name} a échoué: {e}")
        res = {"success": False, "error": str(e)[:300]}
    await _log("executed", {"tool": name, "args": args, "user_id": user_id,
                            "success": bool(res.get("success", True))})
    return res


# ==================== Appel modèle (passerelle + chaîne de repli) ====================
def _now_line():
    jdays = ["lundi", "mardi", "mercredi", "jeudi", "vendredi", "samedi", "dimanche"]
    n = datetime.now(timezone.utc)
    return f"\n\nAUJOURD'HUI : {jdays[n.weekday()]} {n.strftime('%d/%m/%Y')} (ISO {n.strftime('%Y-%m-%d')})."


async def _gemini_call(contents, system):
    """Un tour de génération avec outils. Essaie la chaîne de modèles. Retourne (response, model)."""
    tools = _gemini_tools()
    last_err = None
    for mdl in NEO_MODELS:
        def _call(_m=mdl):
            cfg = _t.GenerateContentConfig(
                system_instruction=system, tools=tools,
                automatic_function_calling=_t.AutomaticFunctionCallingConfig(disable=True),
            )
            return _client.models.generate_content(model=_m, contents=contents, config=cfg)
        try:
            t0 = time.time()
            resp = await asyncio.to_thread(_call)
            await _log("llm", {"model": mdl, "latency_ms": int((time.time() - t0) * 1000)})
            # Gemini renvoie parfois une réponse VIDE (ni texte ni appel d'outil), surtout avec
            # beaucoup d'outils -> on bascule sur le modèle suivant de la chaîne au lieu d'abandonner.
            has_fc = bool(getattr(resp, "function_calls", None))
            try:
                has_txt = bool((resp.text or "").strip())
            except Exception:
                has_txt = False
            if has_fc or has_txt:
                return resp, mdl
            last_err = "empty_response"
            logger.warning(f"neo: modèle {mdl} a renvoyé une réponse vide, essai du modèle suivant")
            continue
        except Exception as e:
            last_err = e
            logger.warning(f"neo: modèle {mdl} a échoué: {e}")
            continue
    raise RuntimeError(f"all_neo_models_failed: {last_err}")


async def _claude_text(system: str, user_text: str) -> str:
    """Appel Claude (Anthropic Messages API, HTTP direct, sans dépendance) pour le jugement stratégique."""
    if not ANTHROPIC_API_KEY:
        raise RuntimeError("no_anthropic_key")
    import requests

    def _call():
        r = requests.post(
            "https://api.anthropic.com/v1/messages",
            headers={"x-api-key": ANTHROPIC_API_KEY, "anthropic-version": "2023-06-01", "content-type": "application/json"},
            json={"model": NEO_STRATEGIC_MODEL, "max_tokens": 1500, "system": system,
                  "messages": [{"role": "user", "content": user_text}]},
            timeout=60)
        r.raise_for_status()
        data = r.json()
        return "".join(b.get("text", "") for b in data.get("content", []) if b.get("type") == "text").strip()
    return await asyncio.to_thread(_call)


async def _gemini_text(system: str, prompt: str) -> str:
    """Génération texte simple (sans outils), chaîne de repli de modèles."""
    if not _client:
        raise RuntimeError("no_gemini")
    last = None
    for mdl in NEO_MODELS:
        def _c(_m=mdl):
            resp = _client.models.generate_content(
                model=_m, contents=prompt,
                config=_t.GenerateContentConfig(system_instruction=system))
            return (getattr(resp, "text", "") or "").strip()
        try:
            t = await asyncio.to_thread(_c)
            if t:
                return t
        except Exception as e:
            last = e
            continue
    raise RuntimeError(f"gemini_text_failed: {last}")


async def _strategic_text(system: str, user_text: str):
    """Lobe stratégique : Claude si dispo (jugement à fort enjeu), sinon repli Gemini. Retourne (texte, modèle)."""
    if ANTHROPIC_API_KEY:
        try:
            t = await _claude_text(system, user_text)
            if t:
                return t, "claude"
        except Exception as e:
            logger.warning(f"neo: Claude indisponible, repli Gemini: {e}")
    return await _gemini_text(system, user_text), "gemini"


def _fr_json(res: dict) -> dict:
    """Réponse d'outil sérialisable pour le modèle (jamais d'ObjectId/datetime brut)."""
    return {"result": json.dumps(res, default=str, ensure_ascii=False)[:4000]}


# ==================== Boucle agentique ====================
def _extract_docx_text(raw: bytes) -> str:
    try:
        import io, docx
        d = docx.Document(io.BytesIO(raw))
        return "\n".join(p.text for p in d.paragraphs if p.text)
    except Exception as e:
        logger.warning(f"docx extract: {e}")
        return ""


def _extract_pptx_text(raw: bytes) -> str:
    try:
        import io
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw))
        out = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for para in shape.text_frame.paragraphs:
                        t = "".join(r.text for r in para.runs)
                        if t:
                            out.append(t)
        return "\n".join(out)
    except Exception as e:
        logger.warning(f"pptx extract: {e}")
        return ""


async def run_neo(messages: list, user_id: str, voice: bool = False, attachments: list = None) -> dict:
    if not _client:
        return {"message": "Néo est momentanément indisponible (clé IA manquante).", "available": False}
    contents = []
    for m in messages:
        role = "model" if m.get("role") == "assistant" else "user"
        contents.append(_t.Content(role=role, parts=[_t.Part.from_text(text=(m.get("content") or "")[:8000])]))

    # Fichiers joints (image/PDF lus nativement par Gemini) -> ajoutés au dernier message utilisateur
    if attachments and contents and contents[-1].role == "user":
        import base64 as _b64
        for att in (attachments or [])[:5]:
            try:
                raw = _b64.b64decode((att.get("data") or "").split(",")[-1])
                mime = att.get("mime_type") or "application/octet-stream"
                name = (att.get("name") or "").lower()
                if raw and (mime.startswith("image/") or mime == "application/pdf"):
                    contents[-1].parts.append(_t.Part.from_bytes(data=raw, mime_type=mime))
                elif raw and (name.endswith(".docx") or "wordprocessingml" in mime):
                    txt = _extract_docx_text(raw)
                    if txt:
                        contents[-1].parts.append(_t.Part.from_text(text=f"[Contenu du fichier {att.get('name')}]\n{txt[:12000]}"))
                elif raw and (name.endswith(".pptx") or "presentationml" in mime):
                    txt = _extract_pptx_text(raw)
                    if txt:
                        contents[-1].parts.append(_t.Part.from_text(text=f"[Contenu du fichier {att.get('name')}]\n{txt[:12000]}"))
            except Exception as e:
                logger.warning(f"neo attachment skip: {e}")

    system = NEO_SYSTEM + _now_line() + await _budget_context() + await _central_memory()
    if voice:
        system += ("\n\n[MODE VOCAL] Tu réponds à l'oral : bref et naturel, 2-3 phrases maximum, "
                   "sans listes à puces, sans markdown, sans emojis. Ton conversationnel, droit au but.")
    pending = []
    actions = []
    for _ in range(MAX_ITERS):
        resp, _mdl = await _gemini_call(contents, system)
        fcs = list(getattr(resp, "function_calls", None) or [])
        if not fcs:
            txt = (getattr(resp, "text", "") or "").strip()
            return {"message": txt or "C'est noté.", "available": True,
                    "pending_actions": pending, "actions_done": actions}
        # rejoue le tour du modèle (avec ses appels d'outils)
        try:
            contents.append(resp.candidates[0].content)
        except Exception:
            contents.append(_t.Content(role="model",
                                       parts=[_t.Part.from_function_call(name=fc.name, args=dict(fc.args or {})) for fc in fcs]))
        tool_parts = []
        for fc in fcs:
            name = fc.name
            args = dict(fc.args or {})
            res = await execute_tool(name, args, user_id)
            if res.get("pending"):
                pending.append({"action_id": res["action_id"], "name": name, "args": args})
            elif res.get("success", True):
                actions.append({"name": name})
            tool_parts.append(_t.Part.from_function_response(name=name, response=_fr_json(res)))
        contents.append(_t.Content(role="tool", parts=tool_parts))

    return {"message": "J'ai atteint la limite d'étapes. Reformule ou demande la suite.",
            "available": True, "pending_actions": pending, "actions_done": actions}


# ==================== Streaming (« il parle pendant qu'il fait ») ====================
# Libellés FR des étapes d'outils affichées EN DIRECT dans le chat : (en cours…, fait).
_TOOL_LABELS = {
    "web_search": ("Recherche sur le web…", "Recherche web terminée"),
    "get_invoice_pdf": ("Récupération du document…", "Document prêt"),
    "delegate_task": ("Délégation à un sous-agent…", "Sous-tâche terminée"),
    "send_to_cowork": ("Envoi à Cowork…", "Tâche envoyée à Cowork"),
    "generate_image": ("Génération du visuel…", "Visuel généré"),
    "search_contacts": ("Recherche de contacts…", "Contacts trouvés"),
    "get_contact": ("Lecture de la fiche…", "Fiche lue"),
    "get_contact_history": ("Lecture de l'historique…", "Historique lu"),
    "read_emails": ("Lecture des emails…", "Emails lus"),
    "list_leads": ("Lecture des leads…", "Leads listés"),
    "get_budget_summary": ("Lecture du budget…", "Budget lu"),
    "get_bank_balance": ("Lecture du solde Qonto…", "Solde bancaire lu"),
    "list_transactions": ("Lecture des transactions…", "Transactions lues"),
    "list_overdue_invoices": ("Recherche des impayés…", "Impayés listés"),
    "list_tasks": ("Lecture des tâches…", "Tâches listées"),
    "get_health_score": ("Calcul du score de santé…", "Score calculé"),
    "strategic_review": ("Analyse stratégique…", "Point stratégique prêt"),
    "activity_report": ("Génération du reporting…", "Reporting prêt"),
    "prep_meeting": ("Préparation du RDV…", "RDV préparé"),
    "enrich_company": ("Recherche sur l'entreprise…", "Entreprise enrichie"),
    "list_documents": ("Lecture des documents…", "Documents listés"),
    "get_document": ("Lecture du document…", "Document lu"),
    "crm_query": ("Lecture du CRM…", "Données lues"),
    "crm_create": ("Création dans le CRM…", "Créé dans le CRM"),
    "crm_update": ("Mise à jour du CRM…", "CRM mis à jour"),
    "crm_delete": ("Préparation de la suppression…", "Suppression à valider"),
    "create_task": ("Création de la tâche…", "Tâche créée"),
    "mark_task_done": ("Clôture de la tâche…", "Tâche terminée"),
    "schedule_followup": ("Programmation de la relance…", "Relance programmée"),
    "create_contact": ("Création de la fiche…", "Fiche contact créée"),
    "set_contact_status": ("Mise à jour du statut…", "Statut mis à jour"),
    "add_contact_note": ("Ajout d'une note…", "Note ajoutée"),
    "update_contact": ("Mise à jour du contact…", "Contact mis à jour"),
    "create_quote": ("Création du devis…", "Devis créé"),
    "draft_followup_email": ("Rédaction de l'email…", "Email préparé"),
    "remember": ("Mémorisation…", "Mémorisé"),
    "recall": ("Consultation de la mémoire…", "Mémoire consultée"),
    "update_objective": ("Mise à jour de l'objectif…", "Objectif enregistré"),
    "log_day": ("Mise à jour du journal…", "Journal mis à jour"),
    "send_followup": ("Préparation de la relance…", "Relance préparée (à valider)"),
    "merge_contacts": ("Préparation de la fusion…", "Fusion préparée (à valider)"),
}


def _tool_label(name: str, kind: str) -> str:
    """Libellé d'étape pour le front. kind: 'start' | 'done' | 'pending' | 'fail'."""
    ing, done = _TOOL_LABELS.get(name, (f"Néo travaille ({name.replace('_', ' ')})…", name.replace("_", " ")))
    if kind == "start":
        return ing
    if kind == "pending":
        return done if "valider" in done else f"{done} — à valider"
    if kind == "fail":
        return f"{done} — échec"
    return done


async def _gemini_stream_turn(contents, system):
    """Un tour de génération EN STREAMING (chaîne de repli de modèles).
    Async-générateur : yield {"type":"text","delta":...} au fil de l'eau, puis un unique
    {"type":"turn_done","fcs":[...],"content":<Content model>} à la fin du tour.
    Le repli vers le modèle suivant n'est possible que TANT QU'aucun token n'a été émis."""
    tools = _gemini_tools()
    last_err = None
    for mdl in NEO_MODELS:
        text_acc = ""
        fcs = []
        started = False
        try:
            t0 = time.time()
            cfg = _t.GenerateContentConfig(
                system_instruction=system, tools=tools,
                automatic_function_calling=_t.AutomaticFunctionCallingConfig(disable=True),
            )
            stream = await _client.aio.models.generate_content_stream(model=mdl, contents=contents, config=cfg)
            async for chunk in stream:
                # appels d'outils (parts complètes, accumulées au fil des chunks)
                try:
                    for fc in (chunk.function_calls or []):
                        fcs.append(fc)
                except Exception:
                    pass
                # fragment de texte
                try:
                    delta = chunk.text or ""
                except Exception:
                    delta = ""
                if delta:
                    started = True
                    text_acc += delta
                    yield {"type": "text", "delta": delta}
            await _log("llm_stream", {"model": mdl, "latency_ms": int((time.time() - t0) * 1000)})
            # Réponse vide (ni texte ni outil) -> on tente le modèle suivant (comme en non-stream)
            if not (text_acc or fcs):
                last_err = "empty_response"
                logger.warning(f"neo stream: modèle {mdl} a renvoyé une réponse vide, essai du suivant")
                continue
            parts = []
            if text_acc:
                parts.append(_t.Part.from_text(text=text_acc))
            for fc in fcs:
                parts.append(_t.Part.from_function_call(name=fc.name, args=dict(fc.args or {})))
            yield {"type": "turn_done", "fcs": fcs, "content": _t.Content(role="model", parts=parts)}
            return
        except Exception as e:
            last_err = e
            logger.warning(f"neo stream: modèle {mdl} a échoué: {e}")
            if started:
                # Du texte a déjà coulé : on ne peut pas rejouer proprement -> on clôt le tour avec l'acquis.
                parts = []
                if text_acc:
                    parts.append(_t.Part.from_text(text=text_acc))
                for fc in fcs:
                    parts.append(_t.Part.from_function_call(name=fc.name, args=dict(fc.args or {})))
                yield {"type": "turn_done", "fcs": fcs, "content": _t.Content(role="model", parts=parts)}
                return
            continue
    raise RuntimeError(f"all_neo_models_failed: {last_err}")


async def run_neo_stream(messages: list, user_id: str, voice: bool = False, attachments: list = None):
    """Version STREAMING de run_neo : même boucle agentique, mais yield les events au fil de l'eau.
    Events : text/tool/pending/done/error. Garde-fous (validation) INCHANGÉS via execute_tool."""
    if not _client:
        yield {"type": "error", "detail": "Néo est momentanément indisponible (clé IA manquante)."}
        return
    contents = []
    for m in messages:
        role = "model" if m.get("role") == "assistant" else "user"
        contents.append(_t.Content(role=role, parts=[_t.Part.from_text(text=(m.get("content") or "")[:8000])]))

    # Fichiers joints (image/PDF lus nativement par Gemini) -> ajoutés au dernier message utilisateur
    if attachments and contents and contents[-1].role == "user":
        import base64 as _b64
        for att in (attachments or [])[:5]:
            try:
                raw = _b64.b64decode((att.get("data") or "").split(",")[-1])
                mime = att.get("mime_type") or "application/octet-stream"
                name = (att.get("name") or "").lower()
                if raw and (mime.startswith("image/") or mime == "application/pdf"):
                    contents[-1].parts.append(_t.Part.from_bytes(data=raw, mime_type=mime))
                elif raw and (name.endswith(".docx") or "wordprocessingml" in mime):
                    txt = _extract_docx_text(raw)
                    if txt:
                        contents[-1].parts.append(_t.Part.from_text(text=f"[Contenu du fichier {att.get('name')}]\n{txt[:12000]}"))
                elif raw and (name.endswith(".pptx") or "presentationml" in mime):
                    txt = _extract_pptx_text(raw)
                    if txt:
                        contents[-1].parts.append(_t.Part.from_text(text=f"[Contenu du fichier {att.get('name')}]\n{txt[:12000]}"))
            except Exception as e:
                logger.warning(f"neo attachment skip: {e}")

    system = NEO_SYSTEM + _now_line() + await _budget_context() + await _central_memory()
    if voice:
        system += ("\n\n[MODE VOCAL] Tu réponds à l'oral : bref et naturel, 2-3 phrases maximum, "
                   "sans listes à puces, sans markdown, sans emojis. Ton conversationnel, droit au but.")
    pending = []
    actions = []
    for _ in range(MAX_ITERS):
        fcs = []
        content_obj = None
        async for ev in _gemini_stream_turn(contents, system):
            if ev["type"] == "text":
                yield ev
            elif ev["type"] == "turn_done":
                fcs = ev["fcs"]
                content_obj = ev["content"]
        if not fcs:
            # Le texte final a déjà été streamé -> on clôt.
            yield {"type": "done", "actions_done": actions, "pending_actions": pending}
            return
        # rejoue le tour du modèle (avec ses appels d'outils)
        contents.append(content_obj)
        tool_parts = []
        for fc in fcs:
            name = fc.name
            args = dict(fc.args or {})
            yield {"type": "tool", "name": name, "phase": "start", "label": _tool_label(name, "start")}
            res = await execute_tool(name, args, user_id)
            if res.get("pending"):
                pending.append({"action_id": res["action_id"], "name": name, "args": args})
                yield {"type": "pending", "action_id": res["action_id"], "name": name, "args": args}
                yield {"type": "tool", "name": name, "phase": "done", "ok": True, "label": _tool_label(name, "pending")}
            elif res.get("success", True):
                actions.append({"name": name})
                yield {"type": "tool", "name": name, "phase": "done", "ok": True, "label": _tool_label(name, "done")}
            else:
                yield {"type": "tool", "name": name, "phase": "done", "ok": False, "label": _tool_label(name, "fail")}
            tool_parts.append(_t.Part.from_function_response(name=name, response=_fr_json(res)))
        contents.append(_t.Content(role="tool", parts=tool_parts))

    # Limite d'étapes atteinte (parité avec run_neo) : on émet un message de clôture.
    yield {"type": "text", "delta": "J'ai atteint la limite d'étapes. Reformule ou demande la suite."}
    yield {"type": "done", "actions_done": actions, "pending_actions": pending}


# ==================== Cerveau Claude (tool-use Anthropic) — moteur alternatif ====================
def _anthropic_tools():
    return [{"name": t["name"], "description": t["description"], "input_schema": t["params"]} for t in TOOLS]


async def _claude_messages_call(system: str, messages: list, tools: list):
    import requests
    def _call():
        r = requests.post(
            "https://api.anthropic.com/v1/messages",
            headers={"x-api-key": ANTHROPIC_API_KEY, "anthropic-version": "2023-06-01", "content-type": "application/json"},
            json={"model": NEO_STRATEGIC_MODEL, "max_tokens": 2000, "system": system, "tools": tools, "messages": messages},
            timeout=70)
        r.raise_for_status()
        return r.json()
    return await asyncio.to_thread(_call)


async def run_neo_claude(messages: list, user_id: str, voice: bool = False) -> dict:
    """Même boucle agentique que run_neo mais pilotée par Claude (tool-use Anthropic)."""
    import json as _json
    if not ANTHROPIC_API_KEY:
        return await run_neo(messages, user_id, voice=voice)
    system = NEO_SYSTEM + _now_line() + await _budget_context() + await _central_memory()
    if voice:
        system += ("\n\n[MODE VOCAL] Tu réponds à l'oral : bref et naturel, 2-3 phrases maximum, "
                   "sans listes à puces, sans markdown, sans emojis. Ton conversationnel, droit au but.")
    conv = [{"role": ("assistant" if m.get("role") == "assistant" else "user"),
             "content": (m.get("content") or "")[:8000]} for m in messages if m.get("content")]
    tools = _anthropic_tools()
    pending, actions = [], []
    for _ in range(MAX_ITERS):
        data = await _claude_messages_call(system, conv, tools)
        blocks = data.get("content", []) or []
        tool_uses = [b for b in blocks if b.get("type") == "tool_use"]
        if not tool_uses:
            txt = "".join(b.get("text", "") for b in blocks if b.get("type") == "text").strip()
            return {"message": txt or "C'est noté.", "available": True,
                    "pending_actions": pending, "actions_done": actions, "brain": "claude"}
        conv.append({"role": "assistant", "content": blocks})
        results = []
        for tu in tool_uses:
            res = await execute_tool(tu.get("name"), tu.get("input") or {}, user_id)
            if res.get("pending"):
                pending.append({"action_id": res["action_id"], "name": tu.get("name"), "args": tu.get("input") or {}})
            elif res.get("success", True):
                actions.append({"name": tu.get("name")})
            results.append({"type": "tool_result", "tool_use_id": tu.get("id"),
                            "content": _json.dumps(res, ensure_ascii=False, default=str)[:6000]})
        conv.append({"role": "user", "content": results})
    return {"message": "J'ai atteint la limite d'étapes. Reformule ou demande la suite.",
            "available": True, "pending_actions": pending, "actions_done": actions, "brain": "claude"}


# ==================== Cerveau Claude EN STREAMING (V2) ====================
async def _claude_stream_turn(system: str, conv: list, tools: list):
    """Un tour Claude EN STREAMING (Anthropic Messages API, stream:true, via httpx async).
    Async-générateur : yield {"type":"text","delta":...} au fil de l'eau, puis un unique
    {"type":"turn_done","blocks":[...],"tool_uses":[...]} (blocks = contenu assistant à rejouer)."""
    import httpx
    import json as _json
    headers = {"x-api-key": ANTHROPIC_API_KEY, "anthropic-version": "2023-06-01", "content-type": "application/json"}
    payload = {"model": NEO_STRATEGIC_MODEL, "max_tokens": 2000, "system": system,
               "tools": tools, "messages": conv, "stream": True}
    blocks = {}   # index -> {"type":"text","text":...} | {"type":"tool_use","id","name","json"}
    order = []
    t0 = time.time()
    async with httpx.AsyncClient(timeout=httpx.Timeout(120.0, connect=15.0)) as client:
        async with client.stream("POST", "https://api.anthropic.com/v1/messages",
                                 headers=headers, json=payload) as resp:
            resp.raise_for_status()
            async for line in resp.aiter_lines():
                if not line or not line.startswith("data:"):
                    continue
                data = line[5:].strip()
                if not data:
                    continue
                try:
                    ev = _json.loads(data)
                except Exception:
                    continue
                etype = ev.get("type")
                if etype == "content_block_start":
                    idx = ev.get("index")
                    cb = ev.get("content_block") or {}
                    if cb.get("type") == "tool_use":
                        blocks[idx] = {"type": "tool_use", "id": cb.get("id"), "name": cb.get("name"), "json": ""}
                    else:
                        blocks[idx] = {"type": "text", "text": ""}
                    order.append(idx)
                elif etype == "content_block_delta":
                    idx = ev.get("index")
                    d = ev.get("delta") or {}
                    b = blocks.get(idx)
                    if not b:
                        continue
                    if d.get("type") == "text_delta":
                        t = d.get("text") or ""
                        if t:
                            b["text"] += t
                            yield {"type": "text", "delta": t}
                    elif d.get("type") == "input_json_delta":
                        b["json"] += d.get("partial_json") or ""
                elif etype == "message_stop":
                    break
    await _log("llm_stream", {"model": "claude", "latency_ms": int((time.time() - t0) * 1000)})
    content, tool_uses = [], []
    for idx in order:
        b = blocks[idx]
        if b["type"] == "text":
            if b["text"].strip():
                content.append({"type": "text", "text": b["text"]})
        else:
            inp = {}
            if (b.get("json") or "").strip():
                try:
                    inp = _json.loads(b["json"])
                except Exception:
                    inp = {}
            content.append({"type": "tool_use", "id": b["id"], "name": b["name"], "input": inp})
            tool_uses.append({"id": b["id"], "name": b["name"], "input": inp})
    yield {"type": "turn_done", "blocks": content, "tool_uses": tool_uses}


async def run_neo_claude_stream(messages: list, user_id: str, voice: bool = False):
    """Version STREAMING de run_neo_claude (V2). Même boucle agentique, events au fil de l'eau.
    Repli sur le stream Gemini si pas de clé Anthropic. Garde-fous (validation) inchangés."""
    import json as _json
    if not ANTHROPIC_API_KEY:
        async for ev in run_neo_stream(messages, user_id, voice=voice):
            yield ev
        return
    system = NEO_SYSTEM + _now_line() + await _budget_context() + await _central_memory()
    if voice:
        system += ("\n\n[MODE VOCAL] Tu réponds à l'oral : bref et naturel, 2-3 phrases maximum, "
                   "sans listes à puces, sans markdown, sans emojis. Ton conversationnel, droit au but.")
    conv = [{"role": ("assistant" if m.get("role") == "assistant" else "user"),
             "content": (m.get("content") or "")[:8000]} for m in messages if m.get("content")]
    tools = _anthropic_tools()
    pending, actions = [], []
    for _ in range(MAX_ITERS):
        tool_uses, blocks = [], None
        async for ev in _claude_stream_turn(system, conv, tools):
            if ev["type"] == "text":
                yield ev
            elif ev["type"] == "turn_done":
                tool_uses = ev["tool_uses"]
                blocks = ev["blocks"]
        if not tool_uses:
            yield {"type": "done", "actions_done": actions, "pending_actions": pending}
            return
        conv.append({"role": "assistant", "content": blocks})
        results = []
        for tu in tool_uses:
            name = tu["name"]
            args = tu.get("input") or {}
            yield {"type": "tool", "name": name, "phase": "start", "label": _tool_label(name, "start")}
            res = await execute_tool(name, args, user_id)
            if res.get("pending"):
                pending.append({"action_id": res["action_id"], "name": name, "args": args})
                yield {"type": "pending", "action_id": res["action_id"], "name": name, "args": args}
                yield {"type": "tool", "name": name, "phase": "done", "ok": True, "label": _tool_label(name, "pending")}
            elif res.get("success", True):
                actions.append({"name": name})
                yield {"type": "tool", "name": name, "phase": "done", "ok": True, "label": _tool_label(name, "done")}
            else:
                yield {"type": "tool", "name": name, "phase": "done", "ok": False, "label": _tool_label(name, "fail")}
            results.append({"type": "tool_result", "tool_use_id": tu["id"],
                            "content": _json.dumps(res, ensure_ascii=False, default=str)[:6000]})
        conv.append({"role": "user", "content": results})
    yield {"type": "text", "delta": "J'ai atteint la limite d'étapes. Reformule ou demande la suite."}
    yield {"type": "done", "actions_done": actions, "pending_actions": pending}


# ==================== Choix AUTOMATIQUE du cerveau (routage hybride) ====================
# Signaux d'une demande COMPLEXE / à fort enjeu -> Claude (raisonnement profond).
_COMPLEX_HINTS = (
    "stratég", "analyse", "analyser", "réfléch", "réflexion", "penses-tu", "ton avis",
    "qu'en penses", "conseil", "recommand", "négoci", "convainc", "argumentaire", "argument",
    "proposition", "offre commerciale", "plan d'", "business plan", "pitch", "rédige", "rédiger",
    "écris-moi", "écris un", "écris une", "idée", "brainstorm", "créati", "complexe", "compliqué",
    "décision", "dois-je", "faut-il", "devrais-je", "pourquoi", "compare", "comparais", "vision",
    "prévision", "anticipe", "scénario", "simul", "optimis", "structure", "challenge", "enjeu",
    "long terme", "rentab", "marge", "positionnement", "diagnostic", "audit", "synthèse", "synthétise",
)
# Démarrages typiques d'un lookup/commande SIMPLE -> Gemini (rapide).
_SIMPLE_HINTS = (
    "liste", "montre", "affiche", "combien", "quel", "quelle", "quels", "qui ", "donne-moi",
    "c'est quoi", "statut", "trouve", "cherche", "ajoute", "crée une tâche", "crée la tâche",
    "marque", "trésorerie", "solde", "factures", "leads", "contacts", "rdv", "rendez-vous",
)


def _last_user_text(messages: list) -> str:
    for m in reversed(messages or []):
        if m.get("role") == "user" and (m.get("content") or "").strip():
            return m["content"].strip()
    return ""


async def _classify_brain(text: str) -> str:
    """Mini-classifieur (modèle rapide) pour les cas AMBIGUS : 'simple' -> gemini, 'complexe' -> claude."""
    if not _client:
        return "gemini"
    prompt = ("Tu classes la demande d'un dirigeant à son associé IA. Réponds UNIQUEMENT par 'simple' ou 'complexe'. "
              "'complexe' = stratégie, analyse, jugement, conseil important, rédaction soignée, négociation, décision à enjeu. "
              "'simple' = consultation du CRM, info factuelle, petite action. "
              f"Demande : « {text[:500]} »\nRéponse :")

    def _call():
        cfg = _t.GenerateContentConfig(max_output_tokens=5, temperature=0)
        return _client.models.generate_content(model="gemini-2.5-flash-lite", contents=prompt, config=cfg)
    try:
        resp = await asyncio.wait_for(asyncio.to_thread(_call), timeout=4.0)
        out = (getattr(resp, "text", "") or "").strip().lower()
        return "claude" if "complex" in out else "gemini"
    except Exception as e:
        logger.warning(f"neo classify_brain KO: {e}")
        return "gemini"


async def _resolve_brain(messages: list, attachments=None) -> str:
    """Routage HYBRIDE : heuristique instantanée pour les cas évidents, mini-classifieur pour les ambigus.
    Retourne 'gemini' (rapide, défaut) ou 'claude' (raisonnement profond)."""
    if attachments:
        return "gemini"  # multimodal -> Gemini (Claude n'a pas les pièces jointes ici)
    text = _last_user_text(messages)
    low = text.lower()
    n = len(text)
    if any(h in low for h in _COMPLEX_HINTS):
        return "claude"                       # signal de complexité/enjeu -> Claude
    if n > 320:
        return "claude"                       # demande longue/élaborée
    if n <= 60:
        return "gemini"                        # très court -> lookup/commande
    if any(h in low for h in _SIMPLE_HINTS):
        return "gemini"                        # lookup/commande clair
    return await _classify_brain(text)         # cas ambigu -> mini-classifieur rapide


# ==================== Endpoints ====================
class NeoChatRequest(BaseModel):
    messages: List[dict]
    conversation_id: Optional[str] = None
    mode: Optional[str] = None  # "voice" => réponses orales concises
    attachments: Optional[List[dict]] = None  # [{name, mime_type, data(base64)}] image/PDF
    brain: Optional[str] = None  # "gemini" (défaut) | "claude"


class ConfirmRequest(BaseModel):
    action_id: str


class FeedbackRequest(BaseModel):
    rating: str  # 'up' | 'down'
    message: Optional[str] = None  # le message de Néo évalué
    note: Optional[str] = None     # la correction / le commentaire de Léo
    conversation_id: Optional[str] = None


@router.get("/health")
async def neo_health():
    return {"available": bool(_client), "model_chain": NEO_MODELS, "tools": [t["name"] for t in TOOLS]}


@router.get("/health-score")
async def neo_health_score_endpoint(current_user: dict = Depends(get_current_user)):
    """Score de santé business /100 (déterministe, sans LLM) — pour le cockpit + Néo."""
    return await _compute_health_score()


async def _checkin_payload(moment: str = None) -> dict:
    """Données du check-in (déterministe, chiffres corrigés). Réutilisé par l'endpoint ET le push proactif."""
    h = (datetime.now(timezone.utc).hour - 4) % 24  # Guadeloupe = UTC-4
    if moment not in ("matin", "après-midi", "soir"):
        moment = "matin" if 4 <= h < 12 else ("après-midi" if 12 <= h < 18 else "soir")
    today = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    try:
        checked = bool(await db.neo_memory.find_one({"type": "daily_log", "key": today}))
    except Exception:
        checked = False
    hs = await _compute_health_score()
    bits = []
    if hs.get("overdue_count"):
        bits.append(f"{hs['overdue_count']} facture(s) en retard ({hs['overdue_total']:.0f}€) à relancer")
    if hs.get("hot_leads"):
        bits.append(f"{hs['hot_leads']} lead(s) chaud(s) à traiter")
    if hs.get("pending_devis"):
        bits.append(f"{hs['pending_devis']} devis à valider")
    prio = " · ".join(bits) if bits else "rien d'urgent, avance sur le fond"
    if moment == "soir":
        msg = f"Bonne soirée Léo. Bilan : {prio}. Comment a avancé ta journée ?"
        question = "Raconte-moi où tu en es, je mets le journal à jour."
    else:
        salut = "Bonjour" if moment == "matin" else "Bon après-midi"
        msg = f"{salut} Léo. Aujourd'hui : {prio}. Santé de l'agence : {hs['score']}/100 ({hs['label']})."
        question = None if checked else "Qu'as-tu de prévu aujourd'hui ? Dis-le-moi, je m'organise avec toi."
    return {"moment": moment, "checked_in_today": checked, "score": hs.get("score"),
            "label": hs.get("label"), "message": msg, "question": question, "priorities": bits}


# ==================== Pousse proactive (Phase 3-4 : in-app + WhatsApp) ====================
async def _deposit_notification(ntype: str, title: str, message: str, priority: str = "normal",
                                data: dict = None, dedup_key: str = None) -> bool:
    """Dépose une notif in-app (db.notifications) + push WebSocket temps réel. Dédup optionnelle par jour.
    Sûr : jamais bloquant."""
    try:
        today = datetime.now(timezone.utc).strftime("%Y-%m-%d")
        if dedup_key:
            exists = await db.notifications.find_one({"type": ntype, "data.dedup": dedup_key, "data.day": today})
            if exists:
                return False
        notif = {"id": str(uuid.uuid4()), "type": ntype, "title": title, "message": message,
                 "data": {**(data or {}), "dedup": dedup_key, "day": today},
                 "priority": priority, "target_users": ["all"], "read_by": [],
                 "created_at": datetime.now(timezone.utc).isoformat()}
        await db.notifications.insert_one(notif)
        try:
            from .notifications import manager  # import lazy (évite la circularité au chargement)
            await manager.broadcast(notif)
        except Exception as e:
            logger.warning(f"neo notif broadcast KO: {e}")
        return True
    except Exception as e:
        logger.warning(f"neo deposit notif KO: {e}")
        return False


async def _whatsapp_push(message: str) -> bool:
    """Pousse un message sur le WhatsApp de l'admin via le microservice existant (si admin_phone configuré)."""
    try:
        cfg = await db.settings.find_one({"type": "whatsapp_config"}, {"_id": 0})
        phone = (cfg or {}).get("admin_phone")
        if not phone:
            return False
        url = os.environ.get("WHATSAPP_SERVICE_URL", "http://localhost:3001")
        import requests
        def _call():
            return requests.post(f"{url}/send", json={"phone": phone, "message": message}, timeout=30)
        r = await asyncio.to_thread(_call)
        return getattr(r, "status_code", 0) == 200
    except Exception as e:
        logger.warning(f"neo whatsapp push KO: {e}")
        return False


async def neo_proactive_push(moment: str = None) -> dict:
    """Briefing proactif de Néo (matin/soir) : dépose une notif in-app + pousse WhatsApp si configuré.
    Réutilise le check-in (chiffres corrects), dédupliqué par jour+moment. Sûr : jamais bloquant.
    Appelé par le scheduler (cron) ou manuellement via POST /neo/run-proactive."""
    try:
        payload = await _checkin_payload(moment)
    except Exception as e:
        logger.warning(f"neo proactive payload KO: {e}")
        return {"success": False, "error": str(e)[:160]}
    moment = payload["moment"]
    prio = payload.get("priorities") or []
    title = "Briefing du matin" if moment == "matin" else ("Récap du soir" if moment == "soir" else "Point de Néo")
    level = "high" if any(("retard" in b or "chaud" in b) for b in prio) else "normal"
    deposited = await _deposit_notification("neo_briefing", title, payload["message"], priority=level,
                                            data={"priorities": prio, "score": payload.get("score")},
                                            dedup_key=f"brief-{moment}")
    wa = await _whatsapp_push(f"{title}\n\n{payload['message']}")
    await _log("proactive", {"moment": moment, "inapp": deposited, "whatsapp": wa})
    return {"success": True, "moment": moment, "inapp_deposited": deposited,
            "whatsapp_sent": wa, "message": payload["message"]}


@router.get("/checkin")
async def neo_checkin(current_user: dict = Depends(get_current_user)):
    """Check-in proactif matin/soir (Phase 3) : Néo te briefe à l'ouverture, ton adapté au
    moment de la journée, alertes finances incluses, et pose la bonne question. Déterministe (rapide)."""
    return await _checkin_payload()


@router.post("/run-proactive")
async def neo_run_proactive(moment: Optional[str] = None, current_user: dict = Depends(get_current_user)):
    """Déclenche le briefing proactif de Néo (in-app + WhatsApp). Pour le scheduler (cron) ou test manuel."""
    return await neo_proactive_push(moment)


@router.get("/strategy")
async def neo_strategy(current_user: dict = Depends(get_current_user)):
    """Point stratégique de Néo (Phase 5) — jugement routé vers Claude (repli Gemini)."""
    return await _strategic_review()


@router.get("/treasury")
async def neo_treasury(current_user: dict = Depends(get_current_user)):
    """Tableau de trésorerie Qonto (Phase 7) : soldes par compte + transactions récentes."""
    return await _qonto_treasury(40)


class TtsRequest(BaseModel):
    text: str
    voice_id: Optional[str] = None
    model: Optional[str] = None  # ex: "eleven_flash_v2_5" pour le mode vocal (latence ~minimale)


def _clean_for_speech(text: str) -> str:
    """Nettoie le texte pour la voix : retire emojis et markdown (sinon ElevenLabs lit « astérisque »...)."""
    import re
    text = re.sub(r"[\U0001F000-\U0001FAFF\U00002600-\U000027BF\U0001F1E6-\U0001F1FF️]", "", text)  # emojis / pictos
    text = re.sub(r"[`*_#>|]", "", text)            # markdown
    text = re.sub(r"^\s*[-•]\s*", "", text, flags=re.MULTILINE)  # puces
    text = re.sub(r"\n{2,}", ". ", text)            # paragraphes -> pause
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()


@router.post("/tts")
async def neo_tts(req: TtsRequest, current_user: dict = Depends(get_current_user)):
    """Voix de Néo : texte -> audio/mpeg via ElevenLabs (synthèse vocale premium)."""
    if not ELEVENLABS_API_KEY:
        raise HTTPException(status_code=503, detail="Voix non configurée (ELEVENLABS_API_KEY manquante).")
    text = _clean_for_speech((req.text or "").strip())[:5000]
    if not text:
        raise HTTPException(status_code=400, detail="Texte requis")
    voice_id = (req.voice_id or ELEVENLABS_VOICE_ID).strip()
    import requests
    try:
        r = await asyncio.to_thread(
            requests.post,
            f"https://api.elevenlabs.io/v1/text-to-speech/{voice_id}",
            headers={"xi-api-key": ELEVENLABS_API_KEY, "Content-Type": "application/json", "Accept": "audio/mpeg"},
            json={"text": text, "model_id": (req.model or ELEVENLABS_MODEL),
                  "voice_settings": {"stability": 0.5, "similarity_boost": 0.75, "style": 0.0, "use_speaker_boost": True}},
            timeout=45,
        )
    except Exception as e:
        logger.warning(f"neo TTS injoignable: {e}")
        raise HTTPException(status_code=502, detail=f"Voix injoignable: {str(e)[:160]}")
    if r.status_code != 200:
        body = (r.text or "")[:300]
        logger.warning(f"neo TTS {r.status_code}: {body}")
        raise HTTPException(status_code=502, detail=f"ElevenLabs {r.status_code}: {body}")
    return Response(content=r.content, media_type="audio/mpeg", headers={"Cache-Control": "no-store"})


@router.get("/voices")
async def neo_voices(current_user: dict = Depends(get_current_user)):
    """Liste les voix ElevenLabs disponibles (pour choisir/changer la voix de Néo)."""
    if not ELEVENLABS_API_KEY:
        return {"connected": False, "voices": []}
    import requests
    try:
        r = await asyncio.to_thread(requests.get, "https://api.elevenlabs.io/v1/voices",
                                    headers={"xi-api-key": ELEVENLABS_API_KEY}, timeout=20)
        r.raise_for_status()
        voices = [{"voice_id": v.get("voice_id"), "name": v.get("name"),
                   "labels": v.get("labels", {}), "preview_url": v.get("preview_url")}
                  for v in (r.json().get("voices", []) or [])]
        return {"connected": True, "current": ELEVENLABS_VOICE_ID, "model": ELEVENLABS_MODEL, "voices": voices}
    except Exception as e:
        logger.warning(f"neo voices error: {e}")
        return {"connected": False, "voices": [], "error": str(e)[:160]}


@router.post("/chat")
async def neo_chat(req: NeoChatRequest, current_user: dict = Depends(get_current_user)):
    user_id = current_user.get("user_id") or current_user.get("id")
    msgs = req.messages or []
    if not msgs:
        raise HTTPException(status_code=400, detail="Message requis")
    try:
        brain = (req.brain or "auto").lower()
        if brain not in ("gemini", "claude"):
            brain = await _resolve_brain(msgs, req.attachments)  # choix auto (hybride)
        if brain == "claude" and not req.attachments:
            try:
                result = await run_neo_claude(msgs, user_id, voice=(req.mode == "voice"))
            except Exception as e:
                logger.warning(f"neo claude KO, repli Gemini: {e}")
                result = await run_neo(msgs, user_id, voice=(req.mode == "voice"), attachments=req.attachments)
        else:
            result = await run_neo(msgs, user_id, voice=(req.mode == "voice"), attachments=req.attachments)
    except Exception as e:
        logger.error(f"neo chat error: {e}")
        raise HTTPException(status_code=500, detail=f"Erreur Néo: {str(e)[:200]}")
    conv_id = req.conversation_id or str(uuid.uuid4())
    result["conversation_id"] = conv_id
    try:
        await _persist_conversation(conv_id, user_id, msgs, result.get("message") or "")
    except Exception as e:
        logger.warning(f"neo conv persist: {e}")
    return result


def _sse(obj: dict) -> str:
    """Sérialise un event en ligne SSE (Server-Sent Events) : 'data: {json}\\n\\n'."""
    return "data: " + json.dumps(obj, ensure_ascii=False, default=str) + "\n\n"


async def _persist_safe(conv_id: str, user_id: str, msgs: list, reply: str):
    """Persistance best-effort de la conversation après un flux (jamais bloquant)."""
    try:
        if reply:
            await _persist_conversation(conv_id, user_id, msgs, reply)
    except Exception as e:
        logger.warning(f"neo conv persist (stream): {e}")


@router.post("/chat/stream")
async def neo_chat_stream(req: NeoChatRequest, current_user: dict = Depends(get_current_user)):
    """Chat de Néo EN STREAMING (SSE) : le texte s'écrit au fil de l'eau + étapes d'outils en direct.
    /neo/chat reste le fallback non-stream. Garde-fous (validation humaine) strictement identiques."""
    user_id = current_user.get("user_id") or current_user.get("id")
    msgs = req.messages or []
    if not msgs:
        raise HTTPException(status_code=400, detail="Message requis")
    conv_id = req.conversation_id or str(uuid.uuid4())
    voice = (req.mode == "voice")
    brain = (req.brain or "auto").lower()
    auto = brain not in ("gemini", "claude")

    async def _gen():
        # Choix AUTO du cerveau (hybride) résolu AVANT le 1er event -> on l'annonce dans meta.
        resolved = brain
        if auto:
            try:
                resolved = await _resolve_brain(msgs, req.attachments)
            except Exception as e:
                logger.warning(f"neo resolve_brain KO: {e}")
                resolved = "gemini"
        yield _sse({"type": "meta", "conversation_id": conv_id, "brain": resolved, "auto": auto})
        text_acc = []
        # V2 : Claude est streamé nativement. Avec pièces jointes -> Gemini (multimodal).
        use_claude = (resolved == "claude" and not req.attachments)
        gen = (run_neo_claude_stream(msgs, user_id, voice=voice) if use_claude
               else run_neo_stream(msgs, user_id, voice=voice, attachments=req.attachments))
        try:
            async for ev in gen:
                if ev.get("type") == "text":
                    text_acc.append(ev.get("delta") or "")
                yield _sse(ev)
        except Exception as e:
            logger.error(f"neo chat stream error: {e}")
            # Repli serveur : si Claude casse AVANT tout texte, on bascule sur le flux Gemini.
            if use_claude and not text_acc:
                logger.warning("neo: repli Claude -> Gemini stream")
                try:
                    async for ev in run_neo_stream(msgs, user_id, voice=voice):
                        if ev.get("type") == "text":
                            text_acc.append(ev.get("delta") or "")
                        yield _sse(ev)
                except Exception as e2:
                    logger.error(f"neo chat stream (repli Gemini) KO: {e2}")
                    yield _sse({"type": "error", "detail": str(e2)[:200]})
            else:
                yield _sse({"type": "error", "detail": str(e)[:200]})
        await _persist_safe(conv_id, user_id, msgs, "".join(text_acc))

    return StreamingResponse(_gen(), media_type="text/event-stream",
                             headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no",
                                      "Connection": "keep-alive"})


def _conv_title(msgs: list) -> str:
    """Titre de conversation = 1er message utilisateur, tronqué."""
    for m in msgs:
        if (m.get("role") == "user") and (m.get("content") or "").strip():
            t = m["content"].strip().splitlines()[0]
            return (t[:50] + "…") if len(t) > 50 else t
    return "Conversation"


async def _persist_conversation(conv_id: str, user_id: str, msgs: list, reply: str):
    """Sauvegarde la conversation (upsert) pour l'historique + reprise."""
    full = [{"role": m.get("role"), "content": m.get("content") or ""} for m in msgs if m.get("content")]
    full.append({"role": "assistant", "content": reply})
    now = datetime.now(timezone.utc).isoformat()
    await db.neo_conversations.update_one(
        {"id": conv_id},
        {"$set": {"messages": full, "title": _conv_title(msgs), "updated_at": now, "user_id": user_id},
         "$setOnInsert": {"id": conv_id, "created_at": now}},
        upsert=True,
    )


@router.get("/conversations")
async def neo_conversations(current_user: dict = Depends(get_current_user)):
    """Liste des conversations de l'utilisateur (récentes d'abord) pour l'historique."""
    uid = current_user.get("user_id") or current_user.get("id")
    rows = await db.neo_conversations.find({"user_id": uid},
        {"_id": 0, "id": 1, "title": 1, "updated_at": 1, "messages": 1}).sort("updated_at", -1).to_list(80)
    out = []
    for r in rows:
        msgs = r.get("messages") or []
        last = (msgs[-1]["content"] if msgs else "") or ""
        out.append({"id": r["id"], "title": r.get("title") or "Conversation",
                    "updated_at": r.get("updated_at"), "count": len(msgs),
                    "preview": (last[:90] + "…") if len(last) > 90 else last})
    return {"conversations": out}


@router.get("/conversations/{conv_id}")
async def neo_conversation_get(conv_id: str, current_user: dict = Depends(get_current_user)):
    """Récupère une conversation complète (pour la reprendre)."""
    uid = current_user.get("user_id") or current_user.get("id")
    r = await db.neo_conversations.find_one({"id": conv_id, "user_id": uid}, {"_id": 0})
    if not r:
        raise HTTPException(status_code=404, detail="Conversation introuvable")
    return r


@router.delete("/conversations/{conv_id}")
async def neo_conversation_delete(conv_id: str, current_user: dict = Depends(get_current_user)):
    uid = current_user.get("user_id") or current_user.get("id")
    await db.neo_conversations.delete_one({"id": conv_id, "user_id": uid})
    return {"success": True}


class CoworkDoneRequest(BaseModel):
    result: Optional[str] = None


@router.get("/cowork-inbox")
async def neo_cowork_inbox(current_user: dict = Depends(get_current_user)):
    """Tâches déposées par Néo pour Cowork (récupérées par le Claude PC via MCP get_cowork_tasks)."""
    rows = await db.cowork_inbox.find({"status": "pending"}, {"_id": 0}).sort("created_at", 1).to_list(50)
    return {"tasks": rows}


@router.post("/cowork-inbox/{task_id}/done")
async def neo_cowork_done(task_id: str, req: CoworkDoneRequest, current_user: dict = Depends(get_current_user)):
    await db.cowork_inbox.update_one({"id": task_id},
        {"$set": {"status": "done", "result": (req.result or "")[:8000], "done_at": datetime.now(timezone.utc).isoformat()}})
    return {"success": True}


@router.get("/image/{image_id}")
async def neo_image(image_id: str):
    """Sert une image générée par Néo (public, id UUID non devinable)."""
    import base64 as _b64
    doc = await db.neo_images.find_one({"id": image_id}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Image introuvable")
    return Response(content=_b64.b64decode(doc["data"]), media_type=doc.get("mime", "image/png"),
                    headers={"Cache-Control": "public, max-age=86400"})


@router.post("/confirm-action")
async def neo_confirm(req: ConfirmRequest, current_user: dict = Depends(get_current_user)):
    user_id = current_user.get("user_id") or current_user.get("id")
    pending = await db.neo_pending_actions.find_one({"id": req.action_id, "status": "pending"})
    if not pending:
        raise HTTPException(status_code=404, detail="Action introuvable ou déjà traitée")
    res = await execute_tool(pending["name"], pending.get("args") or {}, user_id, confirmed=True)
    await db.neo_pending_actions.update_one({"id": req.action_id},
                                            {"$set": {"status": "done", "executed_at": datetime.now(timezone.utc).isoformat()}})
    return {"success": res.get("success", True), "message": res.get("message") or "Action exécutée.", "result": res}


@router.post("/cancel-action")
async def neo_cancel(req: ConfirmRequest, current_user: dict = Depends(get_current_user)):
    await db.neo_pending_actions.update_one({"id": req.action_id}, {"$set": {"status": "cancelled"}})
    return {"success": True, "message": "Action annulée."}


@router.post("/feedback")
async def neo_feedback(req: FeedbackRequest, current_user: dict = Depends(get_current_user)):
    """Apprentissage de Néo (cf. demande Léo « les deux ») :
    - 👍/👎 sur les réponses ; un 👎 avec note (ou un 👍 avec note) devient une LEÇON injectée
      dans le cerveau de Néo à chaque appel ; tout est journalisé comme signal d'apprentissage."""
    uid = current_user.get("user_id") or current_user.get("id")
    now = datetime.now(timezone.utc).isoformat()
    rating = (req.rating or "").strip().lower()
    note = (req.note or "").strip()
    lesson = None
    if rating == "down":
        lesson = "À ÉVITER : " + (note or "réponse jugée insatisfaisante par Léo (reformuler, être plus utile)")
    elif rating == "up" and note:
        lesson = "BONNE APPROCHE (à reproduire) : " + note
    if lesson:
        await db.neo_memory.insert_one({"id": str(uuid.uuid4()), "type": "lesson", "content": lesson[:1000],
                                        "created_at": now, "user_id": uid})
    await db.neo_feedback.insert_one({"id": str(uuid.uuid4()), "rating": rating, "note": note[:1000],
                                      "message": (req.message or "")[:2000], "conversation_id": req.conversation_id,
                                      "user_id": uid, "created_at": now})
    await _log("feedback", {"rating": rating, "has_note": bool(note), "user_id": uid})
    return {"success": True,
            "message": "Merci, je retiens." if rating == "up" else "Compris, je corrige et je le retiens pour la prochaine fois."}
